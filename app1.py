"""
驭长风供应链 — 新员工培训系统 (Streamlit)
Features: account management, training modules, exams, PDF certificates, course upload
"""

import streamlit as st
import sqlite3, hashlib, json, os, time, io, tempfile, re, random
from datetime import datetime
from pathlib import Path
from difflib import SequenceMatcher
import mammoth
import openpyxl
from pptx import Presentation
from pptx.util import Pt

# --------------- DB ---------------
DB_PATH = os.path.join(os.path.dirname(__file__), "training.db")


def get_db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    conn = get_db()
    c = conn.cursor()
    c.executescript("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        display_name TEXT NOT NULL,
        emp_id TEXT DEFAULT '',
        department TEXT DEFAULT '',
        role TEXT DEFAULT 'user',
        created_at TEXT DEFAULT (datetime('now'))
    );
    CREATE TABLE IF NOT EXISTS progress (
        user_id INTEGER,
        module_id INTEGER,
        chapter_idx INTEGER,
        read_done INTEGER DEFAULT 0,
        PRIMARY KEY (user_id, module_id, chapter_idx)
    );
    CREATE TABLE IF NOT EXISTS exam_results (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        module_id INTEGER,
        score INTEGER,
        answers TEXT,
        taken_at TEXT DEFAULT (datetime('now'))
    );
    CREATE TABLE IF NOT EXISTS training_time (
        user_id INTEGER,
        module_id INTEGER,
        seconds INTEGER DEFAULT 0,
        PRIMARY KEY (user_id, module_id)
    );
    CREATE TABLE IF NOT EXISTS modules (
        id INTEGER PRIMARY KEY,
        title TEXT NOT NULL,
        chapters TEXT NOT NULL
    );
    CREATE TABLE IF NOT EXISTS exam_questions (
        module_id INTEGER PRIMARY KEY,
        questions TEXT NOT NULL,
        exam_count INTEGER DEFAULT 0
    );
    """)
    # Seed admin account if not exists
    admin_hash = hashlib.sha256("admin123".encode()).hexdigest()
    c.execute("INSERT OR IGNORE INTO users (username, password_hash, display_name, role) VALUES (?,?,?,?)",
              ("admin", admin_hash, "管理员", "admin"))
    # Seed default modules and exams if empty
    if c.execute("SELECT COUNT(*) FROM modules").fetchone()[0] == 0:
        defaults = get_default_modules()
        for mid, mod in defaults.items():
            c.execute("INSERT INTO modules (id, title, chapters) VALUES (?,?,?)",
                      (mid, mod["title"], json.dumps(mod["chapters"], ensure_ascii=False)))
        default_exams = get_default_exams()
        for mid, qs in default_exams.items():
            c.execute("INSERT INTO exam_questions (module_id, questions) VALUES (?,?)",
                      (mid, json.dumps(qs, ensure_ascii=False)))
    conn.commit()
    conn.close()


# --------------- AUTH HELPERS ---------------
def hash_pw(pw: str) -> str:
    return hashlib.sha256(pw.encode()).hexdigest()


def authenticate(username: str, password: str):
    conn = get_db()
    row = conn.execute("SELECT * FROM users WHERE username=? AND password_hash=?",
                       (username, hash_pw(password))).fetchone()
    conn.close()
    return dict(row) if row else None


def create_user(username, password, display_name, emp_id, department, role="user"):
    conn = get_db()
    try:
        conn.execute("INSERT INTO users (username,password_hash,display_name,emp_id,department,role) VALUES (?,?,?,?,?,?)",
                     (username, hash_pw(password), display_name, emp_id, department, role))
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()


def get_all_users():
    conn = get_db()
    rows = conn.execute("SELECT * FROM users ORDER BY id").fetchall()
    conn.close()
    return [dict(r) for r in rows]


def delete_user(uid):
    conn = get_db()
    conn.execute("DELETE FROM users WHERE id=?", (uid,))
    conn.execute("DELETE FROM progress WHERE user_id=?", (uid,))
    conn.execute("DELETE FROM exam_results WHERE user_id=?", (uid,))
    conn.execute("DELETE FROM training_time WHERE user_id=?", (uid,))
    conn.commit()
    conn.close()


def update_user(uid, **kwargs):
    conn = get_db()
    for k, v in kwargs.items():
        conn.execute(f"UPDATE users SET {k}=? WHERE id=?", (v, uid))
    conn.commit()
    conn.close()


# --------------- PROGRESS DB ---------------
def get_read_checks(user_id, module_id, chapter_count):
    conn = get_db()
    rows = conn.execute("SELECT chapter_idx, read_done FROM progress WHERE user_id=? AND module_id=?",
                        (user_id, module_id)).fetchall()
    conn.close()
    checks = [False] * chapter_count
    for r in rows:
        if r["chapter_idx"] < chapter_count:
            checks[r["chapter_idx"]] = bool(r["read_done"])
    return checks


def mark_chapter_read(user_id, module_id, chapter_idx):
    conn = get_db()
    conn.execute("INSERT OR REPLACE INTO progress (user_id,module_id,chapter_idx,read_done) VALUES (?,?,?,1)",
                 (user_id, module_id, chapter_idx))
    conn.commit()
    conn.close()


def save_exam_result(user_id, module_id, score, answers):
    conn = get_db()
    conn.execute("INSERT INTO exam_results (user_id,module_id,score,answers) VALUES (?,?,?,?)",
                 (user_id, module_id, score, json.dumps(answers)))
    conn.commit()
    conn.close()


def get_best_exam(user_id, module_id):
    conn = get_db()
    row = conn.execute("SELECT MAX(score) as best FROM exam_results WHERE user_id=? AND module_id=?",
                       (user_id, module_id)).fetchone()
    conn.close()
    return row["best"] if row and row["best"] is not None else None


def get_training_time(user_id, module_id):
    conn = get_db()
    row = conn.execute("SELECT seconds FROM training_time WHERE user_id=? AND module_id=?",
                       (user_id, module_id)).fetchone()
    conn.close()
    return row["seconds"] if row else 0


def add_training_time(user_id, module_id, seconds):
    conn = get_db()
    conn.execute("""INSERT INTO training_time (user_id,module_id,seconds) VALUES (?,?,?)
                    ON CONFLICT(user_id,module_id) DO UPDATE SET seconds=seconds+?""",
                 (user_id, module_id, seconds, seconds))
    conn.commit()
    conn.close()
# --------------- DEFAULT TRAINING DATA ---------------
DEPARTMENTS = [
    "运营一部", "运营二部", "运营三部", "运营四部",
    "开发部", "采购部", "财务部", "人力行政部",
    "物控组", "海运物流组", "仓储组", "品质组", "供应链综合组",
]

# Modules stored in session_state for dynamic editing
def get_default_modules():
    return {
        1: {
            "title": "模块一：供应链概述",
            "chapters": [
                {"title": "第一章 为什么要做计划", "html": "<h4>一、计划是什么</h4><p>计划是在了解了市场与竞争对手后，结合自身供应链情况所做出的销售策略。每一个数字背后所代表的一定是一个个的操作。</p><h4>二、计划的三个特征</h4><ul><li><b>所有预测都是错的，但有预测比没有强</b> — 追求及时性，比竞争对手错的少意味着综合成本更低。</li><li><b>预测需要多职能参与</b> — 需要开发、采购/生产、物控、财务甚至总经理参与。</li><li><b>预测需要循环预测，逐渐逼近</b> — 执行中遇到问题需立刻修正计划。</li></ul>"},
                {"title": "第二章 物控的意义", "html": "<h4>一、物控做什么？</h4><ul><li><b>生成需求计划</b>：物控有更多数据及分析能力，达成机会与风险的平衡。</li><li><b>制定发货计划</b>：根据海运周期和工厂生产周期生成发货计划。</li><li><b>与工厂沟通排产</b>：每月沟通未来2-3个月排产及发货调整。</li><li><b>追踪监控销售</b>：每日每周监控销售完成率，及时调整发货和生产。</li><li><b>产品退出机制</b>：每季度根据SKU表现确定品级，升级或淘汰。</li></ul><h4>二、供应链的三道防线</h4><ul><li><b>第一道 — 准确计划</b></li><li><b>第二道 — 安全库存</b>（重点产品3+2模式：海外3个月+工厂2个月）</li><li><b>第三道 — 供应链执行</b>（达成率200%以上全链条救火；70%以下则库存滞销）</li></ul>"},
                {"title": "第三章 供应链流程案例", "html": "<h4>一、销售→需求→生产→采购全链路</h4><p>以春夏季产品为例：1月销售的货，需11月初前发走，9月初前开始备料生产。采购订单从提交到工厂排产一般需7-14天。</p><h4>二、滚动计划</h4><ul><li>对照当前销售完成情况判断后续计划合理性</li><li>确认工厂订单满足发货需求</li><li>生产周期越长，链条越长，反应越慢</li></ul><h4>三、计划注意点</h4><ul><li>没有明显波动不要变更，控制在上下20%浮动</li><li>计划与计划之间要环环相扣</li><li>多变体SKU定价参考库存</li></ul>"},
                {"title": "第四章 产品品级管理", "html": "<h4>品级分类</h4><ul><li><b>爆品</b>：类目TOP，超额备货，变体不超过4个</li><li><b>利润品</b>：毛利率15%以上（美国市场）</li><li><b>新品</b>：上线三个月内，之后定级</li><li><b>清尾品</b>：月销低于30或毛利率低于5%</li><li><b>普通品</b>：以上之外</li></ul>"},
            ],
        },
        2: {
            "title": "模块二：海运发货流程",
            "chapters": [
                {"title": "第一部分 发往海外仓", "html": "<h4>流程</h4><p>每月20号提交滚动五个月销售计划及三个月发货计划。物控审核后反馈数量，运营分配海外仓和官方仓数量后提交。物流部根据计划和库存排柜发货。</p>"},
                {"title": "第二部分 FBA国内直发", "html": "<h4>概述</h4><p>从国内工厂直发FBA仓，2CBM起收。费用=提货报关费+海运费(单价×体积)+税金。</p><h4>操作步骤</h4><ul><li>确认库存 → 严格按审核数量创建货件（超出会被驳回）</li><li>登记直发表 → 发送箱唛条码文件</li><li>箱唛命名：BOX LABEL-SKU-FBA ID-箱数-地址代称</li><li>条码文件至少<b>提前一周</b>给物流部</li></ul><h4>时效</h4><ul><li>整柜到港送仓：7天内</li><li>散货到港送仓：7-12天</li></ul>"},
                {"title": "第三部分 AGL亚马逊物流", "html": "<h4>AGL简介</h4><p>亚马逊官方物流，1CBM起收。优点：上架快、不分仓免锁仓费。缺点：旺季运力紧张。</p><h4>创建货件</h4><ul><li>发货地选国内，勾选亚马逊跨境物流</li><li>比较单点入仓和AMP费用，选更便宜的</li><li>货好时间至少多预留3个工作日</li><li>贸易术语：工厂交货；付款货币：美金</li></ul>"},
                {"title": "第四部分 其他注意事项", "html": "<h4>电池资料</h4><p>带电产品需提供：MSDS、UN38.3检测报告、海运鉴定报告书。</p><h4>欧洲递延资料</h4><ul><li>进口VAT证书</li><li>进口EORI注册证书</li><li>VAT注册公司法人ID及营业执照复印件</li><li>过去三个月完税证明</li><li>进口VAT和平台销售链接</li></ul>"},
            ],
        },
    }


def get_default_exams():
    return {
        1: [
            {"q": "计划的三个特征中，第一个是什么？", "opts": ["所有预测都是错的但有预测比没有强", "预测只需销售部门参与", "计划制定后不需要修改", "预测越乐观越好"], "ans": 0},
            {"q": "供应链第二道防线是？", "opts": ["准确计划", "安全库存", "供应链执行", "价格调整"], "ans": 1},
            {"q": "重点产品备货模式是？", "opts": ["海外1+工厂1", "海外2+工厂3", "海外3+工厂2", "海外5个月"], "ans": 2},
            {"q": "物控核心职能不包括？", "opts": ["生成需求计划", "制定发货计划", "产品广告投放", "产品退出机制"], "ans": 2},
            {"q": "计划达成率低于多少属于库存滞销？", "opts": ["90%", "80%", "70%", "60%"], "ans": 2},
            {"q": "滚动计划中完成率正负多少以内视为合理？", "opts": ["10%", "15%", "20%", "30%"], "ans": 2},
            {"q": "以60天生产+45天海运，1月销售的货最晚何时发走？", "opts": ["前一年10月", "前一年11月初", "前一年12月", "当年1月初"], "ans": 1},
            {"q": "清尾品标准是？", "opts": ["月销<100或毛利<10%", "月销<50或毛利<8%", "月销<30或毛利<5%", "月销<20或毛利<3%"], "ans": 2},
            {"q": "爆品变体数量最多不超过？", "opts": ["2个", "3个", "4个", "6个"], "ans": 2},
            {"q": "采购订单从提交到工厂排产一般需要？", "opts": ["3-5天", "7-14天", "15-20天", "21-30天"], "ans": 1},
        ],
        2: [
            {"q": "每月几号提交滚动销售计划？", "opts": ["1号", "10号", "15号", "20号"], "ans": 3},
            {"q": "FBA国内直发最低按多少CBM收费？", "opts": ["1CBM", "2CBM", "3CBM", "5CBM"], "ans": 1},
            {"q": "AGL最低按多少CBM收费？", "opts": ["0.5CBM", "1CBM", "2CBM", "3CBM"], "ans": 1},
            {"q": "AGL优点不包括？", "opts": ["官方物流有保障", "上架时效快", "不分仓免锁仓费", "旺季运力充足"], "ans": 3},
            {"q": "直发条码文件至少提前多久给物流部？", "opts": ["3天", "5天", "一周", "两周"], "ans": 2},
            {"q": "箱唛打印格式要求？", "opts": ["A4纸", "热敏4×6英寸", "热敏6×8英寸", "A5纸"], "ans": 1},
            {"q": "整柜到港送仓正常时效？", "opts": ["3天", "5天", "7天内", "14天"], "ans": 2},
            {"q": "带电产品发货不需要提供的资料？", "opts": ["MSDS", "UN38.3报告", "海运鉴定报告", "使用说明书"], "ans": 3},
            {"q": "AGL贸易术语应选择？", "opts": ["到岸价", "到门价", "工厂交货", "离岸价"], "ans": 2},
            {"q": "超出审核数量的货件物流部如何处理？", "opts": ["正常发货", "延迟发货", "驳回不发", "加收费用"], "ans": 2},
        ],
    }
# --------------- COURSE PARSER ---------------
def parse_docx_to_chapters(file_bytes: bytes) -> list:
    """Parse DOCX using mammoth with proper heading style mapping."""

    # Map Word heading styles to HTML headings so we can split on them
    style_map = """
    p[style-name='Heading 1'] => h2:fresh
    p[style-name='Heading 2'] => h3:fresh
    p[style-name='Heading 3'] => h4:fresh
    p[style-name='Title'] => h1:fresh
    """
    result = mammoth.convert_to_html(
        io.BytesIO(file_bytes),
        style_map=style_map,
    )
    html = result.value

    # Split by h1/h2 tags (Heading 1 level) into chapters
    # h3 (Heading 2) stays inside the chapter as sub-headings
    chapters = []
    # Use regex to split on <h2> tags (mapped from Heading 1)
    parts = re.split(r'(?=<h[12][^>]*>)', html)

    current_title = None
    current_body = []

    for part in parts:
        part = part.strip()
        if not part:
            continue
        # Check if this part starts with h1 or h2 (chapter-level heading)
        m = re.match(r'<h[12][^>]*>(.*?)</h[12]>', part, re.DOTALL)
        if m:
            # Save previous chapter
            if current_title and current_body:
                chapters.append({"title": current_title, "html": "".join(current_body)})
            current_title = re.sub(r'<[^>]+>', '', m.group(1)).strip()
            # Rest of this part is body content
            body_start = m.end()
            remaining = part[body_start:].strip()
            current_body = [remaining] if remaining else []
        elif current_title:
            current_body.append(part)
        else:
            # Content before any heading — detect bold paragraphs as titles
            text_only = re.sub(r'<[^>]+>', '', part).strip()
            if text_only and len(text_only) < 60:
                current_title = text_only
                current_body = []
            elif text_only:
                current_title = text_only[:50]
                current_body = [part]

    # Flush last chapter
    if current_title and current_body:
        chapters.append({"title": current_title, "html": "".join(current_body)})

    # Fallback: if no chapters found, treat whole doc as one
    if not chapters and html.strip():
        chapters.append({"title": "培训内容", "html": html})

    return chapters


def parse_pptx_to_chapters(file_bytes: bytes) -> list:
    """Parse PPTX using python-pptx, grouping slides into chapters."""

    prs = Presentation(io.BytesIO(file_bytes))
    slides_data = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        slides_data.append(texts)

    # Group slides into chapters
    chapters = []
    current_title = None
    current_html_parts = []

    for paras in slides_data:
        if not paras:
            continue
        title_line = paras[0]
        body_lines = paras[1:]

        # Section divider: short title, little/no body
        if len(body_lines) <= 1 and len(title_line) < 40:
            if current_title and current_html_parts:
                chapters.append({"title": current_title, "html": "".join(current_html_parts)})
                current_html_parts = []
            current_title = title_line
            # If there's one body line, include it
            if body_lines:
                current_html_parts.append(f"<p>{body_lines[0]}</p>")
            continue

        if not current_title:
            current_title = title_line

        html = f"<h4>{title_line}</h4>"
        in_list = False
        for line in body_lines:
            if re.match(r'^[\d]+[.、．]|^[•·\-\*]|^（?\d+）', line):
                if not in_list:
                    html += "<ul>"
                    in_list = True
                clean = re.sub(r'^[\d.、．•·\-\*]+\s*', '', line)
                clean = re.sub(r'^（?\d+）\s*', '', clean)
                html += f"<li>{clean}</li>"
            else:
                if in_list:
                    html += "</ul>"
                    in_list = False
                html += f"<p>{line}</p>"
        if in_list:
            html += "</ul>"
        current_html_parts.append(html)

    if current_title and current_html_parts:
        chapters.append({"title": current_title, "html": "".join(current_html_parts)})

    return chapters


def parse_exam_xlsx(file_bytes: bytes) -> list:
    """Parse Excel file into exam question list.
    Columns: 题目, 选项A, 选项B, 选项C, 选项D, 答案, [题型]
    题型: 单选(default)/多选/简答
    答案: 单选=A, 多选=AB or A,B, 简答=关键词(用|分隔多个)
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
    ws = wb.active
    questions = []
    ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}
    for i, row in enumerate(ws.iter_rows(min_row=2, values_only=True)):
        if not row or not row[0]:
            continue
        q = str(row[0]).strip()
        opts = [str(row[j] or "").strip() for j in range(1, 5)]
        ans_raw = str(row[5] or "").strip() if len(row) > 5 else ""
        qtype_raw = str(row[6] or "").strip() if len(row) > 6 else ""

        # Detect question type
        qtype = "single"
        if qtype_raw in ("多选", "multi", "多选题"):
            qtype = "multi"
        elif qtype_raw in ("简答", "short", "简答题"):
            qtype = "short"
        elif len(ans_raw) > 1 and all(c in "ABCDabcd" for c in ans_raw.replace(",", "").replace("，", "")):
            qtype = "multi"

        if qtype == "short":
            # For short answer, answer is keyword text; options may be empty
            keywords = [k.strip() for k in ans_raw.replace("，", "|").replace(",", "|").split("|") if k.strip()]
            questions.append({"q": q, "type": "short", "keywords": keywords, "answer": ans_raw})
        elif qtype == "multi":
            letters = [c.upper() for c in ans_raw.replace(",", "").replace("，", "") if c.upper() in ans_map]
            ans_indices = sorted(set(ans_map[c] for c in letters))
            valid_opts = [o for o in opts if o]
            if not valid_opts or not ans_indices:
                continue
            questions.append({"q": q, "opts": valid_opts, "ans": ans_indices, "type": "multi"})
        else:
            ans = ans_map.get(ans_raw.upper())
            if any(not o for o in opts) or ans is None:
                continue
            questions.append({"q": q, "opts": opts, "ans": ans, "type": "single"})
    return questions


def parse_exam_docx(file_bytes: bytes) -> list:
    """Parse Word file with question format into exam questions.
    Supports single-choice, multi-choice, and short-answer.
    Multi-choice answer line: 答案: AB or 答案: A,C
    Short-answer: 题目 followed by 答案: text (no ABCD options)
    """
    result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
    lines = [l.strip() for l in result.value.split("\n") if l.strip()]
    questions = []
    i = 0
    opt_pat = re.compile(r'^[A-Da-d][.、．:：\s]')
    ans_line_pat = re.compile(r'(?:答案|answer|正确答案)\s*[：:]\s*(.*)', re.I)
    ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}

    while i < len(lines):
        q_line = lines[i]
        # Try to match as choice question (4 options follow)
        if i + 4 < len(lines):
            cands = lines[i+1:i+5]
            if all(opt_pat.match(c) for c in cands):
                strip_opt = lambda s: re.sub(r'^[A-Da-d][.、．:：\s]\s*', '', s)
                opts = [strip_opt(c) for c in cands]
                # Look for answer line
                if i + 5 < len(lines):
                    ans_match = ans_line_pat.search(lines[i+5])
                    if ans_match:
                        ans_text = ans_match.group(1).strip().upper()
                        letters = [c for c in ans_text.replace(",", "").replace("，", "") if c in ans_map]
                        clean_q = re.sub(r'^[\d]+[.、．:：\s]\s*', '', q_line)
                        if len(letters) > 1:
                            questions.append({"q": clean_q or q_line, "opts": opts,
                                              "ans": sorted(set(ans_map[c] for c in letters)), "type": "multi"})
                        elif len(letters) == 1:
                            questions.append({"q": clean_q or q_line, "opts": opts,
                                              "ans": ans_map[letters[0]], "type": "single"})
                        i += 6
                        continue
                i += 1
                continue

        # Try short answer: question line + answer line (no options)
        if i + 1 < len(lines):
            ans_match = ans_line_pat.search(lines[i+1])
            if ans_match:
                ans_text = ans_match.group(1).strip()
                clean_q = re.sub(r'^[\d]+[.、．:：\s]\s*', '', q_line)
                keywords = [k.strip() for k in ans_text.replace("，", "|").replace(",", "|").split("|") if k.strip()]
                questions.append({"q": clean_q or q_line, "type": "short",
                                  "keywords": keywords, "answer": ans_text})
                i += 2
                continue
        i += 1
    return questions
# --------------- PDF CERTIFICATE ---------------
def generate_certificate_pdf(name, emp_id, dept, module_scores, module_times):
    """Generate a landscape A4 PDF certificate using fpdf2."""
    from fpdf import FPDF
    import math

    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.add_page()
    W, H = 297, 210

    # Gold double border
    pdf.set_draw_color(201, 168, 76)
    pdf.set_line_width(1.5)
    pdf.rect(8, 8, W - 16, H - 16)
    pdf.set_line_width(0.5)
    pdf.rect(12, 12, W - 24, H - 24)

    # Title
    pdf.set_font("Helvetica", "B", 26)
    pdf.set_text_color(26, 60, 110)
    pdf.set_xy(0, 30)
    pdf.cell(W, 10, "TRAINING COMPLETION CERTIFICATE", align="C")

    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(100, 100, 100)
    pdf.set_xy(0, 44)
    pdf.cell(W, 8, "YuChangFeng Supply Chain Training", align="C")

    # Body
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(50, 50, 50)
    pdf.set_xy(0, 62)
    pdf.cell(W, 8, f"Employee: {name}    ID: {emp_id}    Department: {dept}", align="C")
    pdf.set_xy(0, 74)
    pdf.cell(W, 8, "Has successfully completed the New Employee Onboarding Training Program", align="C")
    pdf.set_xy(0, 82)
    pdf.cell(W, 8, "and passed all required examinations.", align="C")

    # Table
    y = 98
    pdf.set_fill_color(26, 60, 110)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "B", 9)
    col_w = [80, 40, 40, 37]
    headers = ["Module", "Score", "Duration", "Result"]
    x_start = (W - sum(col_w)) / 2
    x = x_start
    for i, h in enumerate(headers):
        pdf.set_xy(x, y)
        pdf.cell(col_w[i], 8, h, border=1, fill=True, align="C")
        x += col_w[i]

    pdf.set_text_color(50, 50, 50)
    pdf.set_font("Helvetica", "", 9)
    y += 8
    for mod_name, score, mins in module_scores:
        x = x_start
        vals = [mod_name, str(score), f"{mins} min", "PASS"]
        for i, v in enumerate(vals):
            pdf.set_xy(x, y)
            pdf.cell(col_w[i], 8, v, border=1, align="C")
            x += col_w[i]
        y += 8

    # Seal
    pdf.set_draw_color(192, 57, 43)
    pdf.set_line_width(0.8)
    cx, cy = W / 2, 155
    pdf.circle(cx, cy, 14, style="D")
    pdf.set_font("Helvetica", "B", 7)
    pdf.set_text_color(192, 57, 43)
    pdf.set_xy(cx - 20, cy - 4)
    pdf.cell(40, 5, "YUCHANGFENG", align="C")
    pdf.set_xy(cx - 20, cy + 1)
    pdf.cell(40, 5, "OFFICIAL SEAL", align="C")

    # Signature lines
    pdf.set_draw_color(100, 100, 100)
    pdf.set_line_width(0.3)
    pdf.line(50, 182, 110, 182)
    pdf.line(187, 182, 247, 182)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(100, 100, 100)
    pdf.set_xy(50, 184)
    pdf.cell(60, 6, "Training Manager", align="C")
    pdf.set_xy(187, 184)
    pdf.cell(60, 6, "HR Department", align="C")

    # Certificate number
    now = datetime.now()
    cert_no = f"YCF-{now.strftime('%Y-%m%d')}-{emp_id}"
    pdf.set_font("Helvetica", "", 7)
    pdf.set_text_color(150, 150, 150)
    pdf.set_xy(16, H - 14)
    pdf.cell(100, 5, f"Certificate No: {cert_no}")
    pdf.set_xy(W - 116, H - 14)
    pdf.cell(100, 5, f"Date: {now.strftime('%Y-%m-%d')}", align="R")

    return pdf.output()
# --------------- UI HELPERS ---------------
def fmt_time(seconds):
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    return f"{h}小时{m}分{s}秒" if h else f"{m}分{s}秒"


def get_modules():
    """Load modules from DB."""
    conn = get_db()
    rows = conn.execute("SELECT id, title, chapters FROM modules ORDER BY id").fetchall()
    conn.close()
    return {r["id"]: {"title": r["title"], "chapters": json.loads(r["chapters"])} for r in rows}


def save_module(mid, title, chapters):
    """Insert or replace a module in DB."""
    conn = get_db()
    conn.execute("INSERT OR REPLACE INTO modules (id, title, chapters) VALUES (?,?,?)",
                 (mid, title, json.dumps(chapters, ensure_ascii=False)))
    conn.commit()
    conn.close()


def delete_module_db(mid):
    """Delete a module and its related data from DB."""
    conn = get_db()
    conn.execute("DELETE FROM modules WHERE id=?", (mid,))
    conn.execute("DELETE FROM exam_questions WHERE module_id=?", (mid,))
    conn.execute("DELETE FROM progress WHERE module_id=?", (mid,))
    conn.execute("DELETE FROM exam_results WHERE module_id=?", (mid,))
    conn.execute("DELETE FROM training_time WHERE module_id=?", (mid,))
    conn.commit()
    conn.close()


def get_next_module_id():
    conn = get_db()
    row = conn.execute("SELECT MAX(id) as m FROM modules").fetchone()
    conn.close()
    return (row["m"] or 0) + 1


def get_exams():
    """Load exam questions from DB. Returns {mid: {"questions": [...], "exam_count": int}}"""
    conn = get_db()
    rows = conn.execute("SELECT module_id, questions, exam_count FROM exam_questions").fetchall()
    conn.close()
    result = {}
    for r in rows:
        result[r["module_id"]] = {
            "questions": json.loads(r["questions"]),
            "exam_count": r["exam_count"] or 0,
        }
    return result


def save_exam_questions(mid, questions, exam_count=0):
    """Insert or replace exam questions for a module in DB."""
    conn = get_db()
    conn.execute("INSERT OR REPLACE INTO exam_questions (module_id, questions, exam_count) VALUES (?,?,?)",
                 (mid, json.dumps(questions, ensure_ascii=False), exam_count))
    conn.commit()
    conn.close()


def delete_exam_questions_db(mid):
    conn = get_db()
    conn.execute("DELETE FROM exam_questions WHERE module_id=?", (mid,))
    conn.commit()
    conn.close()


# --------------- PAGE: LOGIN ---------------
def page_login():
    st.markdown("## 驭长风供应链 — 新员工培训系统")
    tab_login, tab_reg = st.tabs(["登录", "注册"])
    with tab_login:
        username = st.text_input("用户名", key="login_user")
        password = st.text_input("密码", type="password", key="login_pw")
        if st.button("登录", type="primary"):
            user = authenticate(username, password)
            if user:
                st.session_state.user = user
                st.session_state.page = "dashboard"
                st.rerun()
            else:
                st.error("用户名或密码错误")
    with tab_reg:
        new_user = st.text_input("用户名", key="reg_user")
        new_pw = st.text_input("密码", type="password", key="reg_pw")
        new_name = st.text_input("姓名", key="reg_name")
        new_emp = st.text_input("工号", key="reg_emp")
        new_dept = st.selectbox("部门", DEPARTMENTS, key="reg_dept")
        if st.button("注册"):
            if not new_user or not new_pw or not new_name:
                st.warning("请填写必填项")
            elif create_user(new_user, new_pw, new_name, new_emp, new_dept):
                st.success("注册成功，请登录")
            else:
                st.error("用户名已存在")


# --------------- PAGE: ADMIN ---------------
def page_admin():
    st.markdown("## 账号管理")
    users = get_all_users()

    # Create new user
    with st.expander("➕ 新建账号"):
        c1, c2 = st.columns(2)
        nu = c1.text_input("用户名", key="adm_nu")
        np = c2.text_input("密码", type="password", key="adm_np")
        c3, c4, c5 = st.columns(3)
        nn = c3.text_input("姓名", key="adm_nn")
        ne = c4.text_input("工号", key="adm_ne")
        nd = c5.selectbox("部门", DEPARTMENTS, key="adm_nd")
        nr = st.selectbox("角色", ["user", "admin"], key="adm_nr")
        if st.button("创建"):
            if nu and np and nn:
                if create_user(nu, np, nn, ne, nd, nr):
                    st.success("创建成功")
                    st.rerun()
                else:
                    st.error("用户名已存在")

    # User table
    st.markdown("### 现有账号")
    for u in users:
        cols = st.columns([2, 2, 2, 1, 1, 1])
        cols[0].write(u["username"])
        cols[1].write(u["display_name"])
        cols[2].write(u["department"])
        cols[3].write(u["role"])
        if u["role"] != "admin" or u["id"] != 1:
            if cols[4].button("重置密码", key=f"rst_{u['id']}"):
                update_user(u["id"], password_hash=hash_pw("123456"))
                st.success(f"{u['username']} 密码已重置为 123456")
            if cols[5].button("删除", key=f"del_{u['id']}"):
                delete_user(u["id"])
                st.rerun()


# --------------- PAGE: DASHBOARD ---------------
def page_dashboard():
    user = st.session_state.user
    modules = get_modules()
    exams = get_exams()
    st.markdown(f"## 培训仪表盘")
    st.markdown(f"**{user['display_name']}** | 工号: {user.get('emp_id','')} | 部门: {user.get('department','')}")

    # Overall progress: count modules that have questions as needing exams
    exam_modules = sum(1 for e in exams.values() if e.get("questions"))
    total_items = sum(len(m["chapters"]) for m in modules.values()) + exam_modules
    completed = 0
    for mid, mod in modules.items():
        checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
        completed += sum(checks)
        best = get_best_exam(user["id"], mid)
        if best is not None and best >= 80:
            completed += 1
    pct = int(completed / total_items * 100) if total_items else 0
    st.progress(pct / 100, text=f"总进度: {pct}%")

    # Module cards
    cols = st.columns(min(len(modules), 3))
    for i, (mid, mod) in enumerate(modules.items()):
        col = cols[i % len(cols)]
        with col:
            checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
            read_done = sum(checks)
            total_ch = len(mod["chapters"])
            best = get_best_exam(user["id"], mid)
            t = get_training_time(user["id"], mid)

            status = "已完成" if (read_done == total_ch and best is not None and best >= 80) else ("进行中" if read_done > 0 else "未开始")
            color = {"已完成": "green", "进行中": "orange", "未开始": "gray"}[status]

            st.markdown(f"""
<div style="border:1px solid #ddd;border-radius:8px;padding:16px;margin-bottom:12px;">
<h4 style="margin:0 0 8px">{mod['title']}</h4>
<span style="background:{color};color:white;padding:2px 8px;border-radius:4px;font-size:0.85em">{status}</span>
<p style="margin:8px 0 4px">章节: {read_done}/{total_ch} | 考试: {best if best is not None else '--'}分</p>
<p style="margin:0;font-size:0.85em;color:#888">学习时长: {fmt_time(t)}</p>
</div>""", unsafe_allow_html=True)

            c1, c2 = st.columns(2)
            if c1.button("学习", key=f"learn_{mid}"):
                st.session_state.page = "module"
                st.session_state.current_module = mid
                st.rerun()
            if c2.button("考试", key=f"exam_{mid}"):
                st.session_state.page = "exam"
                st.session_state.current_module = mid
                st.session_state.exam_started = False
                st.rerun()

    # Certificate button
    all_pass = all(
        (get_best_exam(user["id"], mid) or 0) >= 80
        for mid, e in exams.items() if e.get("questions")
    )
    all_read = all(
        all(get_read_checks(user["id"], mid, len(modules[mid]["chapters"])))
        for mid in modules
    )
    if all_pass and all_read:
        st.markdown("---")
        if st.button("🎓 下载培训合格证书", type="primary"):
            st.session_state.page = "certificate"
            st.rerun()


# --------------- SCORING HELPERS ---------------
def score_short_answer(user_ans: str, item: dict) -> float:
    """Score a short answer question. Returns ratio 0.0~1.0.
    Strategy: keyword hit check + fuzzy similarity against reference answer.
    >= 70% match = full marks, 40~70% = half marks, < 40% = 0.
    """
    user_ans = str(user_ans or "").strip()
    if not user_ans:
        return 0.0
    keywords = item.get("keywords", [])
    ref_answer = item.get("answer", "")

    # Keyword hit: count how many keywords appear in user answer
    if keywords:
        hits = sum(1 for kw in keywords if kw.lower() in user_ans.lower())
        kw_ratio = hits / len(keywords)
    else:
        kw_ratio = 0.0

    # Fuzzy similarity against reference answer
    if ref_answer:
        sim = SequenceMatcher(None, user_ans.lower(), ref_answer.lower()).ratio()
    else:
        sim = 0.0

    # Use the better of the two scores
    best = max(kw_ratio, sim)
    if best >= 0.7:
        return 1.0
    elif best >= 0.4:
        return 0.5
    else:
        return 0.0


# --------------- PAGE: MODULE CONTENT ---------------
def page_module():
    user = st.session_state.user
    mid = st.session_state.current_module
    modules = get_modules()
    mod = modules[mid]

    st.markdown(f"## {mod['title']}")
    if st.button("← 返回仪表盘"):
        st.session_state.page = "dashboard"
        st.rerun()

    checks = get_read_checks(user["id"], mid, len(mod["chapters"]))

    for idx, ch in enumerate(mod["chapters"]):
        with st.expander(f"{'✅' if checks[idx] else '📖'} {ch['title']}", expanded=False):
            st.markdown(ch["html"], unsafe_allow_html=True)
            if not checks[idx]:
                if st.button(f"已确认阅读", key=f"read_{mid}_{idx}"):
                    mark_chapter_read(user["id"], mid, idx)
                    st.rerun()
            else:
                st.success("已阅读")


# --------------- PAGE: EXAM ---------------
def page_exam():
    user = st.session_state.user
    mid = st.session_state.current_module
    modules = get_modules()
    exams = get_exams()
    mod = modules[mid]
    exam_data = exams.get(mid, {})
    all_questions = exam_data.get("questions", []) if isinstance(exam_data, dict) else exam_data
    exam_count = exam_data.get("exam_count", 0) if isinstance(exam_data, dict) else 0

    st.markdown(f"## {mod['title']} — 考试")
    if st.button("← 返回仪表盘", key="exam_back"):
        st.session_state.page = "dashboard"
        st.rerun()

    if not all_questions:
        st.warning("本模块暂无考试题目")
        return

    # Check all chapters read
    checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
    if not all(checks):
        st.warning("请先完成所有章节的阅读确认后再参加考试")
        return

    best = get_best_exam(user["id"], mid)
    if best is not None:
        st.info(f"历史最高分: {best}分 {'(已通过)' if best >= 80 else ''}")

    # Exam timer: 30 min
    EXAM_DURATION = 30 * 60
    # Determine actual question count for this exam
    actual_count = exam_count if (exam_count > 0 and exam_count < len(all_questions)) else len(all_questions)

    if not st.session_state.get("exam_started"):
        st.markdown(f"题库共 **{len(all_questions)}** 题，本次考试随机抽取 **{actual_count}** 题")
        st.markdown(f"每题 **{round(100/actual_count)}** 分，80分及格，限时30分钟")
        if st.button("开始考试", type="primary"):
            # Random selection
            if actual_count < len(all_questions):
                selected = random.sample(all_questions, actual_count)
            else:
                selected = list(all_questions)
                random.shuffle(selected)
            st.session_state.exam_questions = selected
            st.session_state.exam_started = True
            st.session_state.exam_start_time = time.time()
            st.session_state.exam_answers = {}
            st.rerun()
        return

    # Use the pre-selected questions for this exam session
    questions = st.session_state.get("exam_questions", all_questions)

    elapsed = time.time() - st.session_state.exam_start_time
    remaining = max(0, EXAM_DURATION - elapsed)
    rm, rs = divmod(int(remaining), 60)
    st.markdown(f"**剩余时间: {rm:02d}:{rs:02d}**")

    if remaining <= 0:
        st.error("考试时间已到，自动提交")
        _submit_exam(user, mid, questions)
        return

    # Questions
    per_q = round(100 / len(questions))
    answers = st.session_state.get("exam_answers", {})
    for qi, item in enumerate(questions):
        qtype = item.get("type", "single")
        type_label = {"single": "单选", "multi": "多选", "short": "简答"}.get(qtype, "")
        st.markdown(f"**{qi+1}. [{type_label}] {item['q']}** ({per_q}分)")

        if qtype == "single":
            sel = st.radio(
                f"q_{mid}_{qi}",
                options=item["opts"],
                index=answers.get(qi),
                key=f"eq_{mid}_{qi}",
                label_visibility="collapsed",
            )
            if sel is not None:
                answers[qi] = item["opts"].index(sel)
        elif qtype == "multi":
            prev = answers.get(qi, [])
            selected = []
            for oi, opt in enumerate(item["opts"]):
                checked = st.checkbox(f"{chr(65+oi)}. {opt}", value=(oi in prev), key=f"eq_{mid}_{qi}_{oi}")
                if checked:
                    selected.append(oi)
            answers[qi] = sorted(selected)
        elif qtype == "short":
            prev = answers.get(qi, "")
            txt = st.text_area("请输入答案", value=prev, key=f"eq_{mid}_{qi}", height=80)
            answers[qi] = txt

    st.session_state.exam_answers = answers

    if st.button("提交考试", type="primary"):
        unanswered = 0
        for qi, item in enumerate(questions):
            qtype = item.get("type", "single")
            a = answers.get(qi)
            if a is None or (qtype == "multi" and len(a) == 0) or (qtype == "short" and not str(a).strip()):
                unanswered += 1
        if unanswered > 0:
            st.warning(f"还有 {unanswered} 题未作答")
        else:
            _submit_exam(user, mid, questions)


def _submit_exam(user, mid, questions):
    answers = st.session_state.get("exam_answers", {})
    per_q = round(100 / len(questions))
    score = 0
    details = []  # Per-question scoring details for analytics
    for qi, item in enumerate(questions):
        qtype = item.get("type", "single")
        a = answers.get(qi)
        q_score = 0
        if qtype == "single":
            if a == item.get("ans"):
                q_score = per_q
        elif qtype == "multi":
            correct = sorted(item.get("ans", []))
            if sorted(a or []) == correct:
                q_score = per_q
        elif qtype == "short":
            ratio = score_short_answer(a, item)
            q_score = round(per_q * ratio)
        score += q_score
        details.append({"q": item["q"], "type": qtype, "earned": q_score, "max": per_q,
                         "user_ans": a, "correct": item.get("ans") or item.get("answer", "")})
    score = min(score, 100)
    elapsed = time.time() - st.session_state.exam_start_time
    add_training_time(user["id"], mid, int(elapsed))
    # Save with details for analytics
    save_exam_result(user["id"], mid, score, {"answers": answers, "details": details})
    st.session_state.exam_started = False
    st.session_state.exam_questions = None
    st.session_state.last_score = score
    st.session_state.last_details = details
    st.session_state.page = "result"
    st.rerun()


# --------------- PAGE: RESULT ---------------
def page_result():
    score = st.session_state.get("last_score", 0)
    mid = st.session_state.get("current_module", 1)
    modules = get_modules()
    mod = modules[mid]
    passed = score >= 80
    details = st.session_state.get("last_details", [])

    st.markdown(f"## {mod['title']} — 考试结果")
    if passed:
        st.success(f"恭喜通过！得分: {score}/100")
    else:
        st.error(f"未通过，得分: {score}/100 (及格线: 80分)")

    # Show per-question breakdown
    if details:
        with st.expander("逐题详情"):
            for i, d in enumerate(details):
                icon = "✅" if d["earned"] == d["max"] else ("⚠️" if d["earned"] > 0 else "❌")
                st.markdown(f"{icon} **{i+1}. [{d['type']}] {d['q']}** — {d['earned']}/{d['max']}分")

    c1, c2 = st.columns(2)
    if c1.button("返回仪表盘"):
        st.session_state.page = "dashboard"
        st.rerun()
    if not passed:
        if c2.button("重新考试"):
            st.session_state.page = "exam"
            st.session_state.exam_started = False
            st.rerun()


# --------------- PAGE: CERTIFICATE ---------------
def page_certificate():
    user = st.session_state.user
    modules = get_modules()
    exams = get_exams()

    st.markdown("## 培训合格证书")

    module_scores = []
    for mid, mod in modules.items():
        best = get_best_exam(user["id"], mid) or 0
        t = get_training_time(user["id"], mid)
        module_scores.append((mod["title"], best, t // 60))

    # Preview
    st.markdown(f"""
<div style="border:3px double #c9a84c;padding:30px;text-align:center;max-width:600px;margin:auto;background:#fffef7">
<h2 style="color:#1a3c6e">培训合格证书</h2>
<p style="color:#888">驭长风供应链 新员工培训</p>
<p>员工: <b>{user['display_name']}</b> &nbsp; 工号: {user.get('emp_id','')} &nbsp; 部门: {user.get('department','')}</p>
<p>已完成全部培训课程并通过考核</p>
<table style="margin:16px auto;border-collapse:collapse">
<tr style="background:#1a3c6e;color:white"><th style="padding:6px 16px">模块</th><th style="padding:6px 16px">成绩</th><th style="padding:6px 16px">时长</th></tr>
{"".join(f'<tr><td style="padding:6px 16px;border:1px solid #ddd">{n}</td><td style="padding:6px 16px;border:1px solid #ddd;text-align:center">{s}分</td><td style="padding:6px 16px;border:1px solid #ddd;text-align:center">{m}分钟</td></tr>' for n,s,m in module_scores)}
</table>
<p style="color:#888;font-size:0.85em">日期: {datetime.now().strftime('%Y-%m-%d')}</p>
</div>""", unsafe_allow_html=True)

    if st.button("下载PDF证书", type="primary"):
        pdf_bytes = generate_certificate_pdf(
            user["display_name"],
            user.get("emp_id", ""),
            user.get("department", ""),
            module_scores,
            {},
        )
        st.download_button(
            "点击下载",
            data=pdf_bytes,
            file_name=f"certificate_{user.get('emp_id','')}.pdf",
            mime="application/pdf",
        )

    if st.button("返回仪表盘"):
        st.session_state.page = "dashboard"
        st.rerun()


# --------------- PAGE: UPLOAD COURSE ---------------
def page_upload_course():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 上传培训课件")
    st.markdown("支持 `.docx` 和 `.pptx` 格式，系统将自动提取章节内容生成新的培训模块。")

    modules = get_modules()

    # --- Section: add new or replace existing ---
    mode = st.radio("操作", ["新建模块", "替换现有模块课件"], horizontal=True, key="course_mode")

    if mode == "替换现有模块课件" and modules:
        mod_opts = {mid: mod["title"] for mid, mod in modules.items()}
        replace_mid = st.selectbox("选择要替换的模块", options=list(mod_opts.keys()),
                                   format_func=lambda x: mod_opts[x], key="replace_mid")
        mod_title = mod_opts[replace_mid]
    else:
        replace_mid = None
        mod_title = st.text_input("模块名称", placeholder="例: 模块三：仓储管理")

    uploaded = st.file_uploader("选择文件", type=["docx", "pptx"])

    btn_label = "解析并替换课件" if replace_mid else "解析并创建模块"
    if uploaded and mod_title and st.button(btn_label, type="primary"):
        file_bytes = uploaded.read()
        try:
            if uploaded.name.endswith(".docx"):
                chapters = parse_docx_to_chapters(file_bytes)
            else:
                chapters = parse_pptx_to_chapters(file_bytes)

            if not chapters:
                st.error("未能从文件中提取到章节内容")
                return

            if replace_mid:
                # Replace existing module chapters, keep exam
                save_module(replace_mid, mod_title, chapters)
                # Clear read progress for this module for all users
                conn = get_db()
                conn.execute("DELETE FROM progress WHERE module_id=?", (replace_mid,))
                conn.commit()
                conn.close()
                st.success(f"已替换模块 [{mod_title}] 的课件，共 {len(chapters)} 章")
            else:
                new_id = get_next_module_id()
                save_module(new_id, mod_title, chapters)
                save_exam_questions(new_id, [])
                st.success(f"已创建模块 [{mod_title}]，共 {len(chapters)} 章")

            for ch in chapters:
                st.markdown(f"- {ch['title']}")
        except Exception as e:
            st.error(f"解析失败: {e}")


# --------------- PAGE: UPLOAD EXAM ---------------
def page_upload_exam():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 题库管理")

    modules = get_modules()
    if not modules:
        st.warning("请先创建培训模块")
        return

    exams = get_exams()
    mod_opts = {mid: mod["title"] for mid, mod in modules.items()}
    selected_mid = st.selectbox("选择模块", options=list(mod_opts.keys()), format_func=lambda x: mod_opts[x])

    exam_data = exams.get(selected_mid, {})
    current_qs = exam_data.get("questions", []) if isinstance(exam_data, dict) else []
    current_count = exam_data.get("exam_count", 0) if isinstance(exam_data, dict) else 0

    st.markdown(f"当前题库: **{len(current_qs)}** 题 | 每次考试抽取: **{current_count if current_count > 0 else '全部'}** 题")

    # --- Exam count setting ---
    st.markdown("---")
    st.markdown("### 考试抽题设置")
    new_count = st.number_input("每次考试随机抽取题数 (0=全部出题)", min_value=0,
                                max_value=max(len(current_qs), 100), value=current_count, key="exam_count_input")
    if new_count != current_count:
        if st.button("保存抽题设置"):
            save_exam_questions(selected_mid, current_qs, new_count)
            st.success(f"已设置每次考试抽取 {new_count if new_count > 0 else '全部'} 题")
            st.rerun()

    # --- File upload import ---
    st.markdown("---")
    st.markdown("### 从文件导入")
    st.markdown("""
支持: **Excel (.xlsx)** 列: 题目、选项A-D、答案、题型(可选) | **Word (.docx)** 题目+选项+答案行
- 单选: 答案=A  多选: 答案=AB,题型=多选  简答: 选项留空,答案=关键词(|分隔),题型=简答
""")
    uploaded = st.file_uploader("选择题库文件", type=["xlsx", "docx"], key="exam_upload")

    if uploaded and st.button("解析并导入题库", type="primary"):
        file_bytes = uploaded.read()
        try:
            if uploaded.name.endswith(".xlsx"):
                questions = parse_exam_xlsx(file_bytes)
            else:
                questions = parse_exam_docx(file_bytes)

            if not questions:
                st.error("未能从文件中提取到有效题目")
                return

            save_exam_questions(selected_mid, questions, new_count)

            singles = sum(1 for q in questions if q.get("type", "single") == "single")
            multis = sum(1 for q in questions if q.get("type") == "multi")
            shorts = sum(1 for q in questions if q.get("type") == "short")
            st.success(f"已导入 {len(questions)} 题 (单选 {singles} / 多选 {multis} / 简答 {shorts})")
            st.rerun()
        except Exception as e:
            st.error(f"解析失败: {e}")

    # --- Manual entry ---
    st.markdown("---")
    st.markdown("### 手工录入题目")
    with st.form("manual_q", clear_on_submit=True):
        qtype = st.selectbox("题型", ["单选", "多选", "简答"], key="mq_type")
        q_text = st.text_input("题目")
        if qtype in ("单选", "多选"):
            c1, c2 = st.columns(2)
            opt_a = c1.text_input("选项A")
            opt_b = c2.text_input("选项B")
            c3, c4 = st.columns(2)
            opt_c = c3.text_input("选项C")
            opt_d = c4.text_input("选项D")
            if qtype == "单选":
                ans_choice = st.selectbox("正确答案", ["A", "B", "C", "D"])
            else:
                ans_multi = st.multiselect("正确答案(可多选)", ["A", "B", "C", "D"])
        else:
            answer_text = st.text_input("参考答案 (关键词用|分隔)")

        submitted = st.form_submit_button("添加到题库")
        if submitted and q_text:
            ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}
            if qtype == "单选":
                opts = [opt_a, opt_b, opt_c, opt_d]
                if all(opts):
                    new_item = {"q": q_text, "opts": opts, "ans": ans_map[ans_choice], "type": "single"}
                else:
                    st.warning("请填写所有选项")
                    new_item = None
            elif qtype == "多选":
                opts = [opt_a, opt_b, opt_c, opt_d]
                valid_opts = [o for o in opts if o]
                if len(valid_opts) >= 2 and ans_multi:
                    new_item = {"q": q_text, "opts": valid_opts,
                                "ans": sorted(ans_map[c] for c in ans_multi), "type": "multi"}
                else:
                    st.warning("至少2个选项和1个答案")
                    new_item = None
            else:
                keywords = [k.strip() for k in answer_text.replace("，", "|").replace(",", "|").split("|") if k.strip()]
                new_item = {"q": q_text, "type": "short", "keywords": keywords, "answer": answer_text}

            if new_item:
                updated_qs = current_qs + [new_item]
                save_exam_questions(selected_mid, updated_qs, new_count)
                st.success(f"已添加，当前共 {len(updated_qs)} 题")
                st.rerun()

    # --- Preview current questions ---
    if current_qs:
        st.markdown("---")
        st.markdown(f"### 当前题库 ({len(current_qs)} 题)")
        for i, item in enumerate(current_qs):
            qtype = item.get("type", "single")
            label = {"single": "单选", "multi": "多选", "short": "简答"}.get(qtype, "")
            cols = st.columns([10, 1])
            cols[0].markdown(f"**{i+1}. [{label}] {item['q']}**")
            if cols[1].button("删除", key=f"delq_{selected_mid}_{i}"):
                updated = current_qs[:i] + current_qs[i+1:]
                save_exam_questions(selected_mid, updated, new_count)
                st.rerun()


# --------------- PAGE: MANAGE COURSES ---------------
def page_manage_courses():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 课件管理")

    modules = get_modules()
    exams = get_exams()

    if not modules:
        st.info("暂无培训模块")
        return

    for mid, mod in list(modules.items()):
        exam_data = exams.get(mid, {})
        qs = exam_data.get("questions", []) if isinstance(exam_data, dict) else []
        with st.expander(f"{mod['title']}  ({len(mod['chapters'])} 章 | {len(qs)} 题)"):
            # Rename
            new_name = st.text_input("模块名称", value=mod["title"], key=f"rename_{mid}")
            if new_name != mod["title"]:
                if st.button("保存名称", key=f"savename_{mid}"):
                    save_module(mid, new_name, mod["chapters"])
                    st.success("已更新")
                    st.rerun()

            # Chapter list
            st.markdown("**章节列表:**")
            for idx, ch in enumerate(mod["chapters"]):
                st.markdown(f"{idx+1}. {ch['title']}")

            # Delete module
            st.markdown("---")
            st.markdown("**危险操作**")
            if st.button(f"删除整个模块", key=f"delmod_{mid}", type="secondary"):
                st.session_state[f"confirm_del_{mid}"] = True

            if st.session_state.get(f"confirm_del_{mid}"):
                st.warning(f"确认删除 [{mod['title']}]？此操作将清除该模块的所有学员进度和考试记录。")
                c1, c2 = st.columns(2)
                if c1.button("确认删除", key=f"yes_del_{mid}", type="primary"):
                    delete_module_db(mid)
                    st.session_state.pop(f"confirm_del_{mid}", None)
                    st.success("已删除")
                    st.rerun()
                if c2.button("取消", key=f"no_del_{mid}"):
                    st.session_state.pop(f"confirm_del_{mid}", None)
                    st.rerun()

            # Delete exam only
            if qs:
                if st.button(f"清空题库 ({len(qs)} 题)", key=f"delexam_{mid}"):
                    save_exam_questions(mid, [])
                    st.success("题库已清空")
                    st.rerun()


# --------------- PAGE: STUDENT ANALYTICS ---------------
def page_analytics():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 考生分析")

    conn = get_db()
    users = [dict(r) for r in conn.execute(
        "SELECT id, username, display_name, department, emp_id FROM users WHERE role='user' ORDER BY id"
    ).fetchall()]

    if not users:
        st.info("暂无普通用户")
        conn.close()
        return

    modules = get_modules()
    exams = get_exams()

    # --- Overview table ---
    st.markdown("### 全员概览")
    header_cols = ["姓名", "部门", "工号"] + [mod["title"] for mod in modules.values()] + ["总进度"]
    rows = []
    for u in users:
        row = [u["display_name"], u["department"], u["emp_id"]]
        total_ch = 0
        done_ch = 0
        total_exams = 0
        passed_exams = 0
        for mid, mod in modules.items():
            ch_count = len(mod["chapters"])
            checks = get_read_checks(u["id"], mid, ch_count)
            read_n = sum(checks)
            best = get_best_exam(u["id"], mid)
            total_ch += ch_count
            done_ch += read_n
            exam_data = exams.get(mid, {})
            has_exam = bool(exam_data.get("questions") if isinstance(exam_data, dict) else exam_data)
            if has_exam:
                total_exams += 1
                if best is not None and best >= 80:
                    passed_exams += 1
            score_str = f"{best}分" if best is not None else "--"
            status = f"{read_n}/{ch_count}章 | {score_str}"
            row.append(status)
        total_items = total_ch + total_exams
        done_items = done_ch + passed_exams
        pct = round(done_items / total_items * 100) if total_items else 0
        row.append(f"{pct}%")
        rows.append(row)

    import pandas as pd
    df = pd.DataFrame(rows, columns=header_cols)
    st.dataframe(df, use_container_width=True, hide_index=True)

    # --- Individual detail ---
    st.markdown("---")
    st.markdown("### 个人详细分析")
    user_opts = {u["id"]: f"{u['display_name']} ({u['emp_id']})" for u in users}
    sel_uid = st.selectbox("选择考生", options=list(user_opts.keys()), format_func=lambda x: user_opts[x])

    sel_user = next(u for u in users if u["id"] == sel_uid)
    st.markdown(f"**{sel_user['display_name']}** | 部门: {sel_user['department']} | 工号: {sel_user['emp_id']}")

    for mid, mod in modules.items():
        st.markdown(f"#### {mod['title']}")
        # Reading progress
        ch_count = len(mod["chapters"])
        checks = get_read_checks(sel_uid, mid, ch_count)
        read_n = sum(checks)
        st.progress(read_n / ch_count if ch_count else 0, text=f"阅读进度: {read_n}/{ch_count}")

        # Exam history
        exam_rows = conn.execute(
            "SELECT score, answers, taken_at FROM exam_results WHERE user_id=? AND module_id=? ORDER BY taken_at DESC",
            (sel_uid, mid)
        ).fetchall()

        if exam_rows:
            best = max(r["score"] for r in exam_rows)
            st.markdown(f"考试次数: **{len(exam_rows)}** | 最高分: **{best}** | 最近: **{exam_rows[0]['score']}**")

            # Show score trend
            scores = [r["score"] for r in reversed(exam_rows)]
            if len(scores) > 1:
                st.line_chart({"得分": scores}, height=150)

            # Weak areas from latest attempt
            latest_ans = exam_rows[0]["answers"]
            try:
                ans_data = json.loads(latest_ans)
                details = ans_data.get("details", []) if isinstance(ans_data, dict) else []
            except (json.JSONDecodeError, TypeError):
                details = []

            if details:
                wrong = [d for d in details if d["earned"] < d["max"]]
                if wrong:
                    with st.expander(f"薄弱环节 ({len(wrong)} 题)"):
                        for d in wrong:
                            st.markdown(f"- **[{d.get('type','single')}]** {d['q']} (得分 {d['earned']}/{d['max']})")
                else:
                    st.success("最近一次考试全部答对")
        else:
            st.caption("尚未参加考试")

        t = get_training_time(sel_uid, mid)
        st.caption(f"学习时长: {fmt_time(t)}")

    conn.close()


# --------------- SIDEBAR & MAIN ---------------
def main():
    st.set_page_config(page_title="驭长风供应链培训系统", page_icon="📚", layout="wide")
    init_db()

    # Init session state
    if "page" not in st.session_state:
        st.session_state.page = "login"

    # Not logged in
    if "user" not in st.session_state:
        page_login()
        return

    user = st.session_state.user

    # Sidebar navigation
    with st.sidebar:
        st.markdown(f"### 👤 {user['display_name']}")
        st.caption(f"{user.get('department','')} | {user.get('emp_id','')}")
        st.markdown("---")

        if st.button("📊 培训仪表盘", use_container_width=True):
            st.session_state.page = "dashboard"
            st.rerun()

        if user["role"] == "admin":
            if st.button("👥 账号管理", use_container_width=True):
                st.session_state.page = "admin"
                st.rerun()

        if user["role"] == "admin":
            if st.button("📁 上传课件", use_container_width=True):
                st.session_state.page = "upload_course"
                st.rerun()

            if st.button("📝 上传题库", use_container_width=True):
                st.session_state.page = "upload_exam"
                st.rerun()

            if st.button("⚙️ 课件管理", use_container_width=True):
                st.session_state.page = "manage_courses"
                st.rerun()

            if st.button("📈 考生分析", use_container_width=True):
                st.session_state.page = "analytics"
                st.rerun()

        st.markdown("---")
        if st.button("退出登录", use_container_width=True):
            for k in list(st.session_state.keys()):
                del st.session_state[k]
            st.rerun()

    # Page routing
    page = st.session_state.page
    if page == "dashboard":
        page_dashboard()
    elif page == "admin" and user["role"] == "admin":
        page_admin()
    elif page == "module":
        page_module()
    elif page == "exam":
        page_exam()
    elif page == "result":
        page_result()
    elif page == "certificate":
        page_certificate()
    elif page == "upload_course":
        page_upload_course()
    elif page == "upload_exam":
        page_upload_exam()
    elif page == "manage_courses":
        page_manage_courses()
    elif page == "analytics":
        page_analytics()
    else:
        page_dashboard()


if __name__ == "__main__":
    main()
