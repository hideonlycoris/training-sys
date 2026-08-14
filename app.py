"""
新员工培训系统 (Streamlit)
Premium Enterprise Edition
Features: account management, training modules, exams, PDF certificates, course upload, analytics
"""


import re
import google.generativeai as genai
import streamlit as st
import json, os, time, io, random
from datetime import datetime
from pathlib import Path
import mammoth
import openpyxl
from pptx import Presentation
from styles import inject_global_css
from styles import render_login, render_dashboard_header
from styles import render_stats_card, render_progress_card, render_module_card_v2
from styles import render_exam_timer_v2, render_question_v2, render_result_v2, render_certificate_v2
from db import (
    init_db, hash_pw, authenticate, create_user, get_all_users,
    delete_user, update_user, get_read_checks, mark_chapter_read,
    save_exam_result, get_best_exam, get_training_time, add_training_time,
    get_modules, save_module, delete_module_db, get_next_module_id,
    get_exams, save_exam_questions, delete_exam_questions_db,
    seed_admin, seed_default_modules, get_supabase,
)

# --- 云端适配：移除 Windows COM 依赖，云端不支持 Office 自动转 PDF ---

# --------------- FILE STORAGE (云端适配) ---------------
import urllib.parse

UPLOAD_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# 云端适配：移除本地服务器和 Office 转 PDF，改为直接下载/内嵌显示


# --------------- FILE STORAGE ---------------
UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# --------------- DB (from db.py) ---------------
# init_db, hash_pw, authenticate, create_user, get_all_users,
# delete_user, update_user, get_read_checks, mark_chapter_read,
# save_exam_result, get_best_exam, get_training_time, add_training_time,
# get_modules, save_module, delete_module_db, get_next_module_id,
# get_exams, save_exam_questions, delete_exam_questions_db,
# seed_admin, seed_default_modules, get_supabase
# are all imported from db.py at the top of this file


# --------------- DEFAULT TRAINING DATA ---------------
# 部门列表从 Secrets 读取（JSON 数组格式），默认为空
def get_departments():
    import json as _json
    dept_str = st.secrets.get("departments", "[]")
    try:
        return _json.loads(dept_str) if isinstance(dept_str, str) else dept_str
    except:
        return []

def get_default_modules():
    “””默认模块为空，由管理员通过界面上传”””
    return {}

def get_default_exams():
    “””默认考题为空，由管理员通过界面上传”””
    return {}

# --------------- COURSE PARSER ---------------
def parse_docx_to_chapters(file_bytes: bytes) -> list:
    """智能 Word 解析器：支持基于编号特征（一、/1./第一部分）的语义切分"""
    
    # 转换 Word 为 HTML
    result = mammoth.convert_to_html(io.BytesIO(file_bytes))
    html = result.value

    # 定义章节标题的识别特征（匹配：一、 1. 第一章 第一部分 等）
    # 同时保留原有的 h1-h3 标签识别
    chapter_regex = re.compile(
        r'(<h[1-3][^>]*>.*?</h[1-3]>|' # 标准标题标签
        r'<p><strong>\s*[一二三四五六七八九十]+[、.].*?</strong></p>|' # 粗体中文编号
        r'<p>\s*[一二三四五六七八九十]+[、.].*?</p>|' # 普通中文编号
        r'<p><strong>\s*\d+[、.].*?</strong></p>|' # 粗体数字编号
        r'<p>\s*\d+[、.].*?</p>|' # 普通数字编号
        r'<p><strong>\s*第[一二三四五六七八九十]+\s*[章节部].*?</strong></p>|' # 第X章
        r'<p>\s*第[一二三四五六七八九十]+\s*[章节部].*?</p>)', # 第X章
        re.IGNORECASE
    )

    parts = chapter_regex.split(html)
    chapters = []
    current_title = "前言/说明"
    current_content = ""

    for part in parts:
        if not part.strip(): continue
        
        # 判断这一段是否匹配标题特征
        if chapter_regex.match(part):
            # 如果之前的章节有实质内容，先保存
            # 过滤掉只有几个字且无后续内容的“假章节”
            if current_content.strip() and len(re.sub(r'<[^>]+>', '', current_content).strip()) > 5:
                chapters.append({"title": current_title, "html": current_content})
            
            # 提取新标题纯文本作为章节名
            new_title = re.sub(r'<[^>]+>', '', part).strip()
            # 清洗掉标题末尾可能的标点
            current_title = re.sub(r'[、.]$', '', new_title)
            current_content = f"<h4>{new_title}</h4>" # 章节内容以标题开始
        else:
            current_content += part

    # 收尾
    if current_content.strip():
        chapters.append({"title": current_title, "html": current_content})

    # 如果解析出来的章节太少且内容很多，尝试按强行平分逻辑（兜底方案）
    if len(chapters) <= 1 and len(html) > 2000:
        # 这里可以加入进一步细分的逻辑，目前先返回整体
        pass

    return chapters

def parse_pptx_to_chapters(file_bytes: bytes) -> list:
    """智能 PPT 解析器：支持过渡页自动合并，确保每个章节都有其实际内容"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    temp_slides = []

    def extract_html(shape):
        res = ""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t: res += f"<p>{t}</p>"
        elif shape.has_table:
            res += "<table border='1' style='border-collapse:collapse; width:100%; font-size:14px;'>"
            for row in shape.table.rows:
                res += "<tr>"
                for cell in row.cells:
                    txt = "<br>".join([p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()])
                    res += f"<td style='padding:4px; border:1px solid #ccc;'>{txt}</td>"
                res += "</tr>"
            res += "</table>"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes: res += extract_html(s)
        return res

    # 1. 提取所有页面的内容
    for slide in prs.slides:
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        content_html = ""
        for shape in slide.shapes:
            if shape == slide.shapes.title: continue
            content_html += extract_html(shape)
        temp_slides.append({"title": title, "content": content_html})

    # 2. 智能逻辑合并
    final_chapters = []
    buffer_title = ""
    buffer_content = ""

    for i, slide in enumerate(temp_slides):
        # 如果本页有标题
        if slide["title"]:
            # 如果之前 buffer 里已经攒了内容，先结算
            if buffer_content.strip():
                final_chapters.append({"title": buffer_title or "模块开始", "html": f"<h4>{buffer_title}</h4>{buffer_content}"})
                buffer_content = ""
            
            # 更新当前的标题 buffer
            buffer_title = slide["title"]
            buffer_content += slide["content"]
        else:
            # 如果本页没标题，说明是上一页的延续内容，直接叠加内容
            if not buffer_title and i == 0: buffer_title = "课程介绍"
            buffer_content += slide["content"]

        # 特殊处理：如果这页结束后 buffer 里的字数太少（说明可能是纯过渡页），
        # 不结算，继续往后攒，直到遇到下一个有内容的页面
    
    # 最后的结算
    if buffer_title or buffer_content:
        # 修正标题，防止只有“PART ONE”这种没意义的标题
        display_title = buffer_title if len(buffer_title) > 2 else "课程章节"
        final_chapters.append({"title": display_title, "html": f"<h4>{buffer_title}</h4>{buffer_content}"})

    # 3. 二次清洗：剔除只有“目录”、“谢谢观看”等无意义章节
    cleaned_chapters = []
    useless_keywords = ["目录", "CONTENTS", "CONTENT", "谢谢", "THANK", "Q&A", "封面"]
    for c in final_chapters:
        text_content = re.sub(r'<[^>]+>', '', c["html"]).strip()
        if any(kw in c["title"].upper() for kw in useless_keywords) and len(text_content) < 50:
            continue
        cleaned_chapters.append(c)

    return cleaned_chapters

def parse_exam_xlsx(file_bytes: bytes) -> list:
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
    ws = wb.active
    questions = []
    ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}
    for i, row in enumerate(ws.iter_rows(min_row=2, values_only=True)):
        if not row or not row[0]:
            continue
        q = str(row[0]).strip()
        opts = [str(row[j] or "").strip() for j in range(1, 5)]
        ans_raw = str(row[5] or "").strip() if len(row) > 5 else ""
        qtype_raw = str(row[6] or "").strip() if len(row) > 6 else ""

        qtype = "single"
        if qtype_raw in ("多选", "multi", "多选题"):
            qtype = "multi"
        elif qtype_raw in ("简答", "short", "简答题"):
            qtype = "short"
        elif len(ans_raw) > 1 and all(c in "ABCDabcd" for c in ans_raw.replace(",", "").replace("，", "")):
            qtype = "multi"

        if qtype == "short":
            keywords = [k.strip() for k in ans_raw.replace("，", "|").replace(",", "|").split("|") if k.strip()]
            questions.append({"q": q, "type": "short", "keywords": keywords, "answer": ans_raw})
        elif qtype == "multi":
            letters = [c.upper() for c in ans_raw.replace(",", "").replace("，", "") if c.upper() in ans_map]
            ans_indices = sorted(set(ans_map[c] for c in letters))
            valid_opts = [o for o in opts if o]
            if not valid_opts or not ans_indices:
                continue
            questions.append({"q": q, "opts": valid_opts, "ans": ans_indices, "type": "multi"})
        else:
            ans = ans_map.get(ans_raw.upper())
            if any(not o for o in opts) or ans is None:
                continue
            questions.append({"q": q, "opts": opts, "ans": ans, "type": "single"})
    return questions

def parse_exam_docx(file_bytes: bytes) -> list:
    result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
    lines = [l.strip() for l in result.value.split("\n") if l.strip()]
    questions = []
    i = 0
    opt_pat = re.compile(r'^[A-Da-d][.、．:：\s]')
    ans_line_pat = re.compile(r'(?:答案|answer|正确答案)\s*[：:]\s*(.*)', re.I)
    ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}

    while i < len(lines):
        q_line = lines[i]
        if i + 4 < len(lines):
            cands = lines[i+1:i+5]
            if all(opt_pat.match(c) for c in cands):
                strip_opt = lambda s: re.sub(r'^[A-Da-d][.、．:：\s]\s*', '', s)
                opts = [strip_opt(c) for c in cands]
                if i + 5 < len(lines):
                    ans_match = ans_line_pat.search(lines[i+5])
                    if ans_match:
                        ans_text = ans_match.group(1).strip().upper()
                        letters = [c for c in ans_text.replace(",", "").replace("，", "") if c in ans_map]
                        clean_q = re.sub(r'^[\d]+[.、．:：\s]\s*', '', q_line)
                        if len(letters) > 1:
                            questions.append({"q": clean_q or q_line, "opts": opts,
                                              "ans": sorted(set(ans_map[c] for c in letters)), "type": "multi"})
                        elif len(letters) == 1:
                            questions.append({"q": clean_q or q_line, "opts": opts,
                                              "ans": ans_map[letters[0]], "type": "single"})
                        i += 6
                        continue
                i += 1
                continue

        if i + 1 < len(lines):
            ans_match = ans_line_pat.search(lines[i+1])
            if ans_match:
                ans_text = ans_match.group(1).strip()
                clean_q = re.sub(r'^[\d]+[.、．:：\s]\s*', '', q_line)
                keywords = [k.strip() for k in ans_text.replace("，", "|").replace(",", "|").split("|") if k.strip()]
                questions.append({"q": clean_q or q_line, "type": "short",
                                  "keywords": keywords, "answer": ans_text})
                i += 2
                continue
        i += 1
    return questions

# --------------- PDF CERTIFICATE ---------------
def generate_certificate_pdf(name, emp_id, dept, module_scores, module_times):
    from fpdf import FPDF

    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.add_page()
    W, H = 297, 210

    pdf.set_draw_color(201, 168, 76)
    pdf.set_line_width(1.5)
    pdf.rect(8, 8, W - 16, H - 16)
    pdf.set_line_width(0.5)
    pdf.rect(12, 12, W - 24, H - 24)

    pdf.set_font("Helvetica", "B", 26)
    pdf.set_text_color(26, 60, 110)
    pdf.set_xy(0, 30)
    pdf.cell(W, 10, "TRAINING COMPLETION CERTIFICATE", align="C")

    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(100, 100, 100)
    pdf.set_xy(0, 44)
    pdf.cell(W, 8, "YuChangFeng Supply Chain Training", align="C")

    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(50, 50, 50)
    pdf.set_xy(0, 62)
    pdf.cell(W, 8, f"Employee: {name}    ID: {emp_id}    Department: {dept}", align="C")
    pdf.set_xy(0, 74)
    pdf.cell(W, 8, "Has successfully completed the New Employee Onboarding Training Program", align="C")
    pdf.set_xy(0, 82)
    pdf.cell(W, 8, "and passed all required examinations.", align="C")

    y = 98
    pdf.set_fill_color(26, 60, 110)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "B", 9)
    col_w = [80, 40, 40, 37]
    headers = ["Module", "Score", "Duration", "Result"]
    x_start = (W - sum(col_w)) / 2
    x = x_start
    for i, h in enumerate(headers):
        pdf.set_xy(x, y)
        pdf.cell(col_w[i], 8, h, border=1, fill=True, align="C")
        x += col_w[i]

    pdf.set_text_color(50, 50, 50)
    pdf.set_font("Helvetica", "", 9)
    y += 8
    for mod_name, score, mins in module_scores:
        x = x_start
        vals = [mod_name, str(score), f"{mins} min", "PASS"]
        for i, v in enumerate(vals):
            pdf.set_xy(x, y)
            pdf.cell(col_w[i], 8, v, border=1, align="C")
            x += col_w[i]
        y += 8

    pdf.set_draw_color(192, 57, 43)
    pdf.set_line_width(0.8)
    cx, cy = W / 2, 155
    pdf.circle(cx, cy, 14, style="D")
    pdf.set_font("Helvetica", "B", 7)
    pdf.set_text_color(192, 57, 43)
    pdf.set_xy(cx - 20, cy - 4)
    pdf.cell(40, 5, "YUCHANGFENG", align="C")
    pdf.set_xy(cx - 20, cy + 1)
    pdf.cell(40, 5, "OFFICIAL SEAL", align="C")

    pdf.set_draw_color(100, 100, 100)
    pdf.set_line_width(0.3)
    pdf.line(50, 182, 110, 182)
    pdf.line(187, 182, 247, 182)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(100, 100, 100)
    pdf.set_xy(50, 184)
    pdf.cell(60, 6, "Training Manager", align="C")
    pdf.set_xy(187, 184)
    pdf.cell(60, 6, "HR Department", align="C")

    now = datetime.now()
    cert_no = f"YCF-{now.strftime('%Y-%m%d')}-{emp_id}"
    pdf.set_font("Helvetica", "", 7)
    pdf.set_text_color(150, 150, 150)
    pdf.set_xy(16, H - 14)
    pdf.cell(100, 5, f"Certificate No: {cert_no}")
    pdf.set_xy(W - 116, H - 14)
    pdf.cell(100, 5, f"Date: {now.strftime('%Y-%m-%d')}", align="R")

    return pdf.output()

# --------------- UI HELPERS ---------------
def fmt_time(seconds):
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    return f"{h}小时{m}分{s}秒" if h else f"{m}分{s}秒"


# --------------- PAGE: LOGIN ---------------
def page_login():
    # Premium login background
    render_login()

    # Form card centered
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("<div style='margin-top:-300px'></div>", unsafe_allow_html=True)
        tab_login, tab_reg = st.tabs(["🔑 Login", "📝 Register"])
        with tab_login:
            username = st.text_input("Username", key="login_user", placeholder="Enter username")
            password = st.text_input("Password", type="password", key="login_pw", placeholder="Enter password")
            if st.button("Sign In", type="primary", use_container_width=True):
                user = authenticate(username, password)
                if user:
                    st.session_state.user = user
                    st.session_state.page = "dashboard"
                    st.rerun()
                else:
                    st.error("Invalid username or password")
        with tab_reg:
            new_user = st.text_input("Username", key="reg_user", placeholder="Set username")
            new_pw = st.text_input("Password", type="password", key="reg_pw", placeholder="Set password")
            new_name = st.text_input("Full Name", key="reg_name", placeholder="Your name")
            new_emp = st.text_input("Employee ID", key="reg_emp", placeholder="Employee ID")
            new_dept = st.selectbox("Department", get_departments(), key="reg_dept")
            if st.button("Create Account", use_container_width=True):
                if not new_user or not new_pw or not new_name:
                    st.warning("Please fill in all required fields")
                elif create_user(new_user, new_pw, new_name, new_emp, new_dept):
                    st.success("Account created! Please sign in.")
                else:
                    st.error("Username already exists")

# --------------- PAGE: ADMIN ---------------
def page_admin():
    st.markdown("## 账号管理")
    users = get_all_users()

    with st.expander("➕ 新建账号"):
        c1, c2 = st.columns(2)
        nu = c1.text_input("用户名", key="adm_nu")
        np = c2.text_input("密码", type="password", key="adm_np")
        c3, c4, c5 = st.columns(3)
        nn = c3.text_input("姓名", key="adm_nn")
        ne = c4.text_input("工号", key="adm_ne")
        nd = c5.selectbox("部门", get_departments(), key="adm_nd")
        nr = st.selectbox("角色", ["user", "admin"], key="adm_nr")
        if st.button("创建"):
            if nu and np and nn:
                if create_user(nu, np, nn, ne, nd, nr):
                    st.success("创建成功")
                    st.rerun()
                else:
                    st.error("用户名已存在")

    st.markdown("### 现有账号")
    for u in users:
        cols = st.columns([2, 2, 2, 1, 1, 1])
        cols[0].write(u["username"])
        cols[1].write(u["display_name"])
        cols[2].write(u["department"])
        cols[3].write(u["role"])
        if u["role"] != "admin" or u["id"] != 1:
            if cols[4].button("重置密码", key=f"rst_{u['id']}"):
                reset_pw = st.secrets.get("default_reset_password", "123456")
                update_user(u["id"], password_hash=hash_pw(reset_pw))
                st.success(f"{u['username']} 密码已重置")
            if cols[5].button("删除", key=f"del_{u['id']}"):
                delete_user(u["id"])
                st.rerun()

# --------------- PAGE: DASHBOARD ---------------
def page_dashboard():
    user = st.session_state.user
    modules = get_modules()
    exams = get_exams()

    # Premium dashboard header
    render_dashboard_header(user['display_name'], user.get('department',''), user.get('emp_id',''))

    # Calculate stats
    exam_modules = sum(1 for e in exams.values() if e.get("questions"))
    total_chapters = sum(len(m["chapters"]) for m in modules.values())
    total_modules = len(modules)
    total_exams = exam_modules

    # Stats cards row
    st.markdown("<div style='height:24px'></div>", unsafe_allow_html=True)
    cols = st.columns(4)
    with cols[0]:
        render_stats_card("📚", "TOTAL MODULES", str(total_modules), "#1a3c6e")
    with cols[1]:
        render_stats_card("📖", "CHAPTERS", str(total_chapters), "#2a5ca8")
    with cols[2]:
        render_stats_card("📝", "EXAMS", str(total_exams), "#c9a84c")
    with cols[3]:
        render_stats_card("⏱", "HOURS", "12+", "#059669")

    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)

    # Progress card
    completed = 0
    for mid, mod in modules.items():
        checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
        completed += sum(checks)
        best = get_best_exam(user["id"], mid)
        if best is not None and best >= 80:
            completed += 1
    pct = int(completed / (total_chapters + exam_modules) * 100) if (total_chapters + exam_modules) else 0
    render_progress_card(pct)

    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)

    # Section header
    st.markdown("""
<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:20px">
  <div>
    <h2 style="font-size:1.25rem;font-weight:700;color:#0f172a;margin:0">Training Modules</h2>
    <p style="color:#64748b;font-size:.9rem;margin:4px 0 0">Complete all modules to receive certification</p>
  </div>
</div>
""", unsafe_allow_html=True)

    # Module cards
    for mid, mod in modules.items():
        checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
        read_done = sum(checks)
        total_ch = len(mod["chapters"])
        best = get_best_exam(user["id"], mid)
        t = get_training_time(user["id"], mid)

        status = "已完成" if (read_done == total_ch and best is not None and best >= 80) else ("进行中" if read_done > 0 else "未开始")
        is_completed = status == "已完成"

        render_module_card_v2(
            title=mod["title"],
            status=status,
            chapters=f"{read_done}/{total_ch}",
            exam=f"{best if best is not None else '--'}",
            duration=fmt_time(t),
            module_id=mid,
            is_completed=is_completed
        )

        # Action buttons
        c1, c2 = st.columns(2)
        with c1:
            if st.button("📖 Start Learning", key=f"learn_{mid}", use_container_width=True):
                st.session_state.page = "module"
                st.session_state.current_module = mid
                st.rerun()
        with c2:
            if st.button("📝 Take Exam", key=f"exam_{mid}", use_container_width=True):
                st.session_state.page = "exam"
                st.session_state.current_module = mid
                st.session_state.exam_started = False
                st.rerun()

        st.markdown("<div style='height:16px'></div>", unsafe_allow_html=True)

    # Certificate button
    all_pass = all(
        (get_best_exam(user["id"], mid) or 0) >= 80
        for mid, e in exams.items() if e.get("questions")
    )
    all_read = all(
        all(get_read_checks(user["id"], mid, len(modules[mid]["chapters"])))
        for mid in modules
    )
    if all_pass and all_read:
        st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)
        if st.button("🎓 Get Your Certificate", type="primary", use_container_width=True):
            st.session_state.page = "certificate"
            st.rerun()

# --------------- SCORING HELPERS ---------------
def score_short_answer(user_ans: str, item: dict) -> tuple[float, str]:
    """基于 Google Gemini 的语义化智能打分，并返回详细判分理由"""
    
    user_ans = str(user_ans or "").strip()
    if not user_ans:
        return 0.0, "学员未作答，计0分。"
    
    ref_answer = item.get("answer", "")
    keywords = item.get("keywords", [])

    # 读取 Gemini 配置开关
    USE_AI_SCORING = st.secrets.get("gemini", {}).get("use_ai_scoring", False)
    
    if USE_AI_SCORING and ref_answer:
        try:
            api_key = st.secrets["gemini"]["api_key"]
            model_name = st.secrets["gemini"].get("model_name", "gemini-2.5-flash")
            
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(model_name)
            
            prompt = f"""你是一个严谨的供应链系统新员工培训考试阅卷助手。
            请对比【标准参考答案】与【学员回答】，判断学员是否答到了核心语义。
            
            【标准参考答案】：{ref_answer}
            【学员回答】：{user_ans}
            
            评分规则（满分1.0）：
            - 核心意思完全一致，即使措辞和描述方式不同，给 1.0
            - 表达了部分核心意思，给 0.5
            - 意思完全无关或错误，给 0.0
            
            请严格输出纯 JSON 字符串，不要带 markdown 标记(如 ```json )。格式必须如下：
            {{"score": 0.5, "reason": "这里写30字以内的评分理由，指出答对了什么或欠缺了什么"}}"""

            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.1,
                )
            )
            
            llm_reply = response.text.strip()
            clean_json = llm_reply.replace("```json", "").replace("```", "").strip()
            
            result = json.loads(clean_json)
            score = float(result.get("score", 0.0))
            reason = result.get("reason", "AI打分成功")
            
            return min(max(score, 0.0), 1.0), f"🤖 阅卷助手: {reason}"
                
        except Exception as e:
            print(f"Gemini 阅卷异常，降级为常规匹配: {e}")
            pass 
            
    # ================= 常规匹配算法 (降级方案) =================
    score = 0.0
    if keywords:
        hits = sum(1 for kw in keywords if kw.lower() in user_ans.lower())
        score = hits / len(keywords) if keywords else 0.0
    
    if ref_answer:
        ref_words = set(ref_answer.lower())
        ans_words = set(user_ans.lower())
        overlap = len(ref_words.intersection(ans_words))
        if len(ref_words) > 0:
            sim = overlap / len(ref_words)
            score = max(score, sim)
            
    if score >= 0.6: 
        return 1.0, "⚙️ 系统常规判定：命中多数核心关键词"
    elif score >= 0.3: 
        return 0.5, "⚙️ 系统常规判定：命中部分核心关键词"
    else: 
        return 0.0, "⚙️ 系统常规判定：未命中核心关键词或偏题"

# --------------- PAGE: MODULE CONTENT ---------------
def page_module():
    import urllib.parse
    user = st.session_state.user
    mid = st.session_state.current_module
    modules = get_modules()
    mod = modules[mid]

    # 品牌化模块头部
    st.markdown(f"""
<div class="module-content-header">
  <h2>📖 {mod['title']}</h2>
</div>
""", unsafe_allow_html=True)

    if st.button("← 返回仪表盘"):
        st.session_state.page = "dashboard"
        st.rerun()

    checks = get_read_checks(user["id"], mid, len(mod["chapters"]))

    for idx, ch in enumerate(mod["chapters"]):
        title_text = ch.get("title", "未命名章节")
        status_icon = "✅" if checks[idx] else "📖"

        with st.expander(f"{status_icon} 课件：{title_text}", expanded=not checks[idx]):
            file_path = ch.get("file_path", "")
            file_type = ch.get("file_type", "pdf")
            old_html = ch.get("html", "")

            if file_path:
                if not os.path.exists(file_path):
                    st.error("⚠️ 文件路径失效，请管理员重新上传此模块。")
                else:
                    if os.path.getsize(file_path) == 0:
                        st.error("⚠️ 该文件转换失败（大小为0字节），请检查本地 Office 是否卡死。")
                    else:
                        # 云端适配：直接提供下载（无法使用本地服务器预览）
                        with open(file_path, "rb") as f:
                            st.download_button(label=f"📥 下载课件 ({file_type.upper()})", data=f, file_name=f"{title_text}.{file_type}", key=f"dl_{mid}_{idx}")
            elif old_html:
                # 注意：此处渲染管理员上传的 HTML 内容
                # 风险：恶意 DOCX/PPTX 可能包含 JavaScript
                # 缓解：仅管理员可上传，且内容在内部使用
                st.markdown(old_html, unsafe_allow_html=True)
            else:
                st.info("该章节暂无内容。")

            st.markdown("<br>", unsafe_allow_html=True)
            if not checks[idx]:
                if st.button(f"✅ 确认完成学习", key=f"read_{mid}_{idx}", type="primary"):
                    mark_chapter_read(user["id"], mid, idx)
                    st.rerun()
            else:
                st.success("✅ 学习已完成")

# --------------- PAGE: EXAM ---------------
def page_exam():
    user = st.session_state.user
    mid = st.session_state.current_module
    modules = get_modules()
    exams = get_exams()
    mod = modules[mid]
    exam_data = exams.get(mid, {})
    all_questions = exam_data.get("questions", []) if isinstance(exam_data, dict) else exam_data
    exam_count = exam_data.get("exam_count", 0) if isinstance(exam_data, dict) else 0

    # Back button
    if st.button("← Back to Dashboard", key="exam_back"):
        st.session_state.page = "dashboard"
        st.rerun()

    # Module title
    st.markdown(f"""
<div style="margin-bottom:24px">
  <h1 style="font-size:1.5rem;font-weight:700;color:#0f172a;margin:0">{mod['title']}</h1>
  <p style="color:#64748b;font-size:.9rem;margin:4px 0 0">Online Examination</p>
</div>
""", unsafe_allow_html=True)

    if not all_questions:
        st.warning("No questions available for this module")
        return

    checks = get_read_checks(user["id"], mid, len(mod["chapters"]))
    if not all(checks):
        st.warning("Please complete all chapter readings before taking the exam")
        return

    best = get_best_exam(user["id"], mid)
    if best is not None:
        st.info(f"Best score: {best}/100 {'(Passed)' if best >= 80 else ''}")

    EXAM_DURATION = 30 * 60
    actual_count = exam_count if (exam_count > 0 and exam_count < len(all_questions)) else len(all_questions)

    if not st.session_state.get("exam_started"):
        st.markdown(f"""
<div style="background:#fff;border-radius:16px;padding:32px;border:1px solid #e2e8f0;margin-bottom:24px">
  <h3 style="font-size:1.1rem;font-weight:600;color:#0f172a;margin-bottom:16px">Exam Details</h3>
  <div style="display:grid;grid-template-columns:repeat(2,1fr);gap:16px">
    <div style="padding:16px;background:#f8fafc;border-radius:10px">
      <div style="font-size:.75rem;color:#94a3b8;text-transform:uppercase;letter-spacing:.5px">Questions</div>
      <div style="font-size:1.5rem;font-weight:700;color:#0f172a">{actual_count}</div>
    </div>
    <div style="padding:16px;background:#f8fafc;border-radius:10px">
      <div style="font-size:.75rem;color:#94a3b8;text-transform:uppercase;letter-spacing:.5px">Passing Score</div>
      <div style="font-size:1.5rem;font-weight:700;color:#0f172a">80/100</div>
    </div>
    <div style="padding:16px;background:#f8fafc;border-radius:10px">
      <div style="font-size:.75rem;color:#94a3b8;text-transform:uppercase;letter-spacing:.5px">Time Limit</div>
      <div style="font-size:1.5rem;font-weight:700;color:#0f172a">30 min</div>
    </div>
    <div style="padding:16px;background:#f8fafc;border-radius:10px">
      <div style="font-size:.75rem;color:#94a3b8;text-transform:uppercase;letter-spacing:.5px">Points Each</div>
      <div style="font-size:1.5rem;font-weight:700;color:#0f172a">{round(100/actual_count)}</div>
    </div>
  </div>
</div>
""", unsafe_allow_html=True)

        if st.button("Start Exam", type="primary", use_container_width=True):
            if actual_count < len(all_questions):
                selected = random.sample(all_questions, actual_count)
            else:
                selected = list(all_questions)
                random.shuffle(selected)
            st.session_state.exam_questions = selected
            st.session_state.exam_started = True
            st.session_state.exam_start_time = time.time()
            st.session_state.exam_answers = {}
            st.rerun()
        return

    questions = st.session_state.get("exam_questions", all_questions)
    elapsed = time.time() - st.session_state.exam_start_time
    remaining = max(0, EXAM_DURATION - elapsed)
    rm, rs = divmod(int(remaining), 60)

    # Timer
    is_warning = remaining <= 300
    render_exam_timer_v2(rm, rs, is_warning)

    if remaining <= 0:
        st.error("Time's up! Auto-submitting...")
        _submit_exam(user, mid, questions)
        return

    per_q = round(100 / len(questions))
    answers = st.session_state.get("exam_answers", {})

    st.markdown("<div style='height:24px'></div>", unsafe_allow_html=True)

    for qi, item in enumerate(questions):
        qtype = item.get("type", "single")

        # Render question card
        render_question_v2(qi + 1, item["q"], qtype, per_q, len(questions))

        if qtype == "single":
            sel = st.radio(
                f"q_{mid}_{qi}",
                options=item["opts"],
                index=answers.get(qi),
                key=f"eq_{mid}_{qi}",
                label_visibility="collapsed",
            )
            if sel is not None:
                answers[qi] = item["opts"].index(sel)
        elif qtype == "multi":
            prev = answers.get(qi, [])
            selected = []
            for oi, opt in enumerate(item["opts"]):
                checked = st.checkbox(f"{chr(65+oi)}. {opt}", value=(oi in prev), key=f"eq_{mid}_{qi}_{oi}")
                if checked:
                    selected.append(oi)
            answers[qi] = sorted(selected)
        elif qtype == "short":
            prev = answers.get(qi, "")
            txt = st.text_area("Enter your answer", value=prev, key=f"eq_{mid}_{qi}", height=100,
                              placeholder="Type your answer here...")
            answers[qi] = txt

    st.session_state.exam_answers = answers

    # Submit button
    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)
    if st.button("Submit Exam", type="primary", use_container_width=True):
        unanswered = 0
        for qi, item in enumerate(questions):
            qtype = item.get("type", "single")
            a = answers.get(qi)
            if a is None or (qtype == "multi" and len(a) == 0) or (qtype == "short" and not str(a).strip()):
                unanswered += 1
        if unanswered > 0:
            st.warning(f"{unanswered} questions unanswered. Please check before submitting.")
        else:
            _submit_exam(user, mid, questions)

def _submit_exam(user, mid, questions):
    answers = st.session_state.get("exam_answers", {})
    per_q = round(100 / len(questions))
    score = 0
    details = []
    
    for qi, item in enumerate(questions):
        qtype = item.get("type", "single")
        a = answers.get(qi)
        q_score = 0
        ai_reason = "" # 初始化判分理由
        
        if qtype == "single":
            if a == item.get("ans"):
                q_score = per_q
        elif qtype == "multi":
            correct = sorted(item.get("ans", []))
            if sorted(a or []) == correct:
                q_score = per_q
        elif qtype == "short":
            # 接收分数比例和判卷理由
            ratio, ai_reason = score_short_answer(a, item)
            q_score = round(per_q * ratio)
            
        score += q_score
        details.append({
            "q": item["q"], 
            "type": qtype, 
            "earned": q_score, 
            "max": per_q,
            "user_ans": a, 
            "correct": item.get("ans") or item.get("answer", ""),
            "ai_reason": ai_reason # 将理由存入详情
        })
        
    score = min(score, 100)
    elapsed = time.time() - st.session_state.exam_start_time
    add_training_time(user["id"], mid, int(elapsed))
    save_exam_result(user["id"], mid, score, {"answers": answers, "details": details})
    
    st.session_state.exam_started = False
    st.session_state.exam_questions = None
    st.session_state.last_score = score
    st.session_state.last_details = details
    st.session_state.page = "result"
    st.rerun()

# --------------- PAGE: RESULT ---------------
def page_result():
    score = st.session_state.get("last_score", 0)
    mid = st.session_state.get("current_module", 1)
    modules = get_modules()
    mod = modules[mid]
    passed = score >= 80
    details = st.session_state.get("last_details", [])

    # Render result
    render_result_v2(score, passed, mod['title'])

    # Details section
    if details:
        st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)
        with st.expander("📝 Detailed Review", expanded=True):
            for i, d in enumerate(details):
                icon = "✅" if d["earned"] == d["max"] else ("⚠️" if d["earned"] > 0 else "❌")
                type_label = {"single": "Single", "multi": "Multi", "short": "Short"}.get(d["type"], "")

                st.markdown(f"**{icon} {i+1}. [{type_label}] {d['q']}**")

                if d['type'] == "short":
                    st.markdown(f"> **Your Answer**: {d.get('user_ans', 'No answer')}")
                    st.markdown(f"> **Reference**: {d['correct']}")
                    st.markdown(f"**Score**: {d['earned']}/{d['max']}")
                    if d.get("ai_reason"):
                        st.info(d["ai_reason"])
                else:
                    st.markdown(f"**Score**: {d['earned']}/{d['max']}")

                st.divider()

    # Action buttons
    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        if st.button("📊 Dashboard", use_container_width=True):
            st.session_state.page = "dashboard"
            st.rerun()
    with c2:
        if not passed:
            if st.button("🔄 Retake Exam", type="primary", use_container_width=True):
                st.session_state.page = "exam"
                st.session_state.exam_started = False
                st.rerun()

# --------------- PAGE: CERTIFICATE ---------------
def page_certificate():
    user = st.session_state.user
    modules = get_modules()

    module_scores = []
    for mid, mod in modules.items():
        best = get_best_exam(user["id"], mid) or 0
        t = get_training_time(user["id"], mid)
        module_scores.append((mod["title"], best, t // 60))

    date_str = datetime.now().strftime('%Y-%m-%d')

    # Render certificate
    render_certificate_v2(
        name=user["display_name"],
        emp_id=user.get("emp_id", ""),
        dept=user.get("department", ""),
        modules=module_scores,
        date=date_str
    )

    # Download button
    st.markdown("<div style='height:32px'></div>", unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        if st.button("📥 Download PDF", type="primary", use_container_width=True):
            pdf_bytes = generate_certificate_pdf(
                user["display_name"],
                user.get("emp_id", ""),
                user.get("department", ""),
                module_scores,
                {},
            )
            st.download_button(
                "Download",
                data=pdf_bytes,
                file_name=f"certificate_{user.get('emp_id','')}.pdf",
                mime="application/pdf",
                use_container_width=True,
            )
    with c2:
        if st.button("📊 Dashboard", use_container_width=True):
            st.session_state.page = "dashboard"
            st.rerun()

# --------------- PAGE: UPLOAD COURSE ---------------
def page_upload_course():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 上传培训课件 (支持原格式 1:1 在线预览)")
    st.info("系统会自动将上传的 Word/PPT 在后台转换为网页可读格式，保证排版原汁原味。")

    modules = get_modules()
    mode = st.radio("操作", ["新建模块", "替换现有模块课件"], horizontal=True, key="course_mode")

    if mode == "替换现有模块课件" and modules:
        mod_opts = {mid: mod["title"] for mid, mod in modules.items()}
        replace_mid = st.selectbox("选择要替换的模块", options=list(mod_opts.keys()),
                                   format_func=lambda x: mod_opts[x], key="replace_mid")
        mod_title = mod_opts[replace_mid]
    else:
        replace_mid = None
        mod_title = st.text_input("模块名称", placeholder="例: 模块三：仓储管理")

    uploaded_files = st.file_uploader("选择课件文件 (支持多选 PDF/DOCX/PPTX)", type=["pdf", "docx", "pptx", "doc", "ppt"], accept_multiple_files=True)

    btn_label = "上传并替换课件" if replace_mid else "上传并创建模块"
    if uploaded_files and mod_title and st.button(btn_label, type="primary"):
        chapters = []
        try:
            for f in uploaded_files:
                # 安全处理文件名：去除路径分隔符，防止路径穿越
                import re as _re
                clean_name = _re.sub(r'[/\\:*?"<>|]', '_', f.name)
                safe_filename = f"{int(time.time())}_{clean_name}"
                file_path = os.path.join(UPLOAD_DIR, safe_filename)

                # 保存原始文件
                with open(file_path, "wb") as out:
                    out.write(f.read())

                ext = f.name.split('.')[-1].lower()
                # 云端适配：直接存储原始文件，提供下载
                chapters.append({"title": f.name, "file_path": file_path, "file_type": ext})

            if replace_mid:
                save_module(replace_mid, mod_title, chapters)
                # 清除该模块的学习进度
                from db import get_supabase
                sb = get_supabase()
                sb.table("training_progress").delete().eq("module_id", replace_mid).execute()
                st.success(f"已替换模块 [{mod_title}] 的课件。")
            else:
                new_id = get_next_module_id()
                save_module(new_id, mod_title, chapters)
                save_exam_questions(new_id, [])
                st.success(f"已创建模块 [{mod_title}]。")

        except Exception as e:
            st.error(f"处理失败: {e}")

# --------------- PAGE: UPLOAD EXAM ---------------
def page_upload_exam():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 题库管理")

    modules = get_modules()
    if not modules:
        st.warning("请先创建培训模块")
        return

    exams = get_exams()
    mod_opts = {mid: mod["title"] for mid, mod in modules.items()}
    selected_mid = st.selectbox("选择模块", options=list(mod_opts.keys()), format_func=lambda x: mod_opts[x])

    exam_data = exams.get(selected_mid, {})
    current_qs = exam_data.get("questions", []) if isinstance(exam_data, dict) else []
    current_count = exam_data.get("exam_count", 0) if isinstance(exam_data, dict) else 0

    st.markdown(f"当前题库: **{len(current_qs)}** 题 | 每次考试抽取: **{current_count if current_count > 0 else '全部'}** 题")

    st.markdown("---")
    st.markdown("### 考试抽题设置")
    new_count = st.number_input("每次考试随机抽取题数 (0=全部出题)", min_value=0,
                                max_value=max(len(current_qs), 100), value=current_count, key="exam_count_input")
    if new_count != current_count:
        if st.button("保存抽题设置"):
            save_exam_questions(selected_mid, current_qs, new_count)
            st.success(f"已设置每次考试抽取 {new_count if new_count > 0 else '全部'} 题")
            st.rerun()

    # --- Batch entry ---
    st.markdown("---")
    st.markdown("### ⚡ 批量快速录入")
    st.markdown("支持直接粘贴文本，每行一题。格式：`题目 | 选项A | 选项B | 选项C | 选项D | 答案(A/B/C/D)`")
    batch_text = st.text_area("在此粘贴题库文本", height=150, key="batch_upload_text")

    if st.button("批量解析入库", type="primary"):
        if batch_text.strip():
            new_items = []
            ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}
            for line in batch_text.split("\n"):
                line = line.strip()
                if not line: continue
                parts = [p.strip() for p in line.split("|")]
                if len(parts) >= 6:
                    q_text, a, b, c, d, ans_str = parts[:6]
                    # 去除干扰字符，提取纯大写字母答案
                    ans_clean = ans_str.upper().replace('答案', '').replace(':', '').replace('：', '').replace(',', '').replace('，', '').strip()
                    
                    # 识别答案中包含了哪些有效选项字母
                    letters = [char for char in ans_clean if char in ans_map]
                    
                    if len(letters) > 1:
                        # 识别为多选题
                        ans_indices = sorted(list(set([ans_map[char] for char in letters])))
                        new_items.append({
                            "q": q_text, 
                            "opts": [a, b, c, d], 
                            "ans": ans_indices, 
                            "type": "multi"
                        })
                    elif len(letters) == 1:
                        # 识别为单选题
                        new_items.append({
                            "q": q_text, 
                            "opts": [a, b, c, d], 
                            "ans": ans_map[letters[0]], 
                            "type": "single"
                        })
            
            if new_items:
                updated_qs = current_qs + new_items
                save_exam_questions(selected_mid, updated_qs, exam_count=10) # 自动设置为抽取10道题
                
                # 统计一下导入了多少单选，多少多选
                singles = sum(1 for q in new_items if q["type"] == "single")
                multis = sum(1 for q in new_items if q["type"] == "multi")
                
                st.success(f"成功批量录入 {len(new_items)} 道题 (单选 {singles} 题, 多选 {multis} 题)！当前考试已自动设置为随机抽取 10 题。")
                st.rerun()
            else:
                st.error("未能解析出有效题目，请检查格式分隔符是否为 '|' 且答案是否包含A/B/C/D")

    st.markdown("---")
    st.markdown("### 从文件导入")
    st.markdown("""
支持: **Excel (.xlsx)** 列: 题目、选项A-D、答案、题型(可选) | **Word (.docx)** 题目+选项+答案行
- 单选: 答案=A  多选: 答案=AB,题型=多选  简答: 选项留空,答案=关键词(|分隔),题型=简答
""")
    uploaded = st.file_uploader("选择题库文件", type=["xlsx", "docx"], key="exam_upload")

    if uploaded and st.button("解析并导入题库"):
        file_bytes = uploaded.read()
        try:
            if uploaded.name.endswith(".xlsx"):
                questions = parse_exam_xlsx(file_bytes)
            else:
                questions = parse_exam_docx(file_bytes)

            if not questions:
                st.error("未能从文件中提取到有效题目")
                return

            save_exam_questions(selected_mid, current_qs + questions, new_count)

            singles = sum(1 for q in questions if q.get("type", "single") == "single")
            multis = sum(1 for q in questions if q.get("type") == "multi")
            shorts = sum(1 for q in questions if q.get("type") == "short")
            st.success(f"已导入 {len(questions)} 题 (单选 {singles} / 多选 {multis} / 简答 {shorts})")
            st.rerun()
        except Exception as e:
            st.error(f"解析失败: {e}")

    # --- Manual entry ---
    st.markdown("---")
    st.markdown("### 手工录入单道题目")
    
    # 【关键修改】：将题型选择器移动到 form 的外面，使其可以实时触发界面更新
    qtype = st.selectbox("题型", ["单选", "多选", "简答"], key="mq_type")
    
    with st.form("manual_q", clear_on_submit=True):
        q_text = st.text_input("题目")
        if qtype in ("单选", "多选"):
            c1, c2 = st.columns(2)
            opt_a = c1.text_input("选项A")
            opt_b = c2.text_input("选项B")
            c3, c4 = st.columns(2)
            opt_c = c3.text_input("选项C")
            opt_d = c4.text_input("选项D")
            if qtype == "单选":
                ans_choice = st.selectbox("正确答案", ["A", "B", "C", "D"])
            else:
                ans_multi = st.multiselect("正确答案(可多选)", ["A", "B", "C", "D"])
        else:
            answer_text = st.text_input("参考答案 (关键词用|分隔)")

        submitted = st.form_submit_button("添加到题库")
        if submitted and q_text:
            ans_map = {"A": 0, "B": 1, "C": 2, "D": 3}
            if qtype == "单选":
                opts = [opt_a, opt_b, opt_c, opt_d]
                if all(opts):
                    new_item = {"q": q_text, "opts": opts, "ans": ans_map[ans_choice], "type": "single"}
                else:
                    st.warning("请填写所有选项")
                    new_item = None
            elif qtype == "多选":
                opts = [opt_a, opt_b, opt_c, opt_d]
                valid_opts = [o for o in opts if o]
                if len(valid_opts) >= 2 and ans_multi:
                    new_item = {"q": q_text, "opts": valid_opts,
                                "ans": sorted(ans_map[c] for c in ans_multi), "type": "multi"}
                else:
                    st.warning("至少2个选项和1个答案")
                    new_item = None
            else:
                keywords = [k.strip() for k in answer_text.replace("，", "|").replace(",", "|").split("|") if k.strip()]
                new_item = {"q": q_text, "type": "short", "keywords": keywords, "answer": answer_text}

            if new_item:
                updated_qs = current_qs + [new_item]
                save_exam_questions(selected_mid, updated_qs, new_count)
                st.success(f"已添加，当前共 {len(updated_qs)} 题")
                st.rerun()

    if current_qs:
        st.markdown("---")
        st.markdown(f"### 当前题库 ({len(current_qs)} 题)")
        for i, item in enumerate(current_qs):
            qtype = item.get("type", "single")
            label = {"single": "单选", "multi": "多选", "short": "简答"}.get(qtype, "")
            cols = st.columns([10, 1])
            cols[0].markdown(f"**{i+1}. [{label}] {item['q']}**")
            if cols[1].button("删除", key=f"delq_{selected_mid}_{i}"):
                updated = current_qs[:i] + current_qs[i+1:]
                save_exam_questions(selected_mid, updated, new_count)
                st.rerun()

# --------------- PAGE: MANAGE COURSES ---------------
def page_manage_courses():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 课件管理")

    modules = get_modules()
    exams = get_exams()

    if not modules:
        st.info("暂无培训模块")
        return

    for mid, mod in list(modules.items()):
        exam_data = exams.get(mid, {})
        qs = exam_data.get("questions", []) if isinstance(exam_data, dict) else []
        with st.expander(f"{mod['title']}  ({len(mod['chapters'])} 章 | {len(qs)} 题)"):
            new_name = st.text_input("模块名称", value=mod["title"], key=f"rename_{mid}")
            if new_name != mod["title"]:
                if st.button("保存名称", key=f"savename_{mid}"):
                    save_module(mid, new_name, mod["chapters"])
                    st.success("已更新")
                    st.rerun()

            st.markdown("**章节列表:**")
            for idx, ch in enumerate(mod["chapters"]):
                st.markdown(f"{idx+1}. {ch['title']}")

            st.markdown("---")
            st.markdown("**危险操作**")
            if st.button(f"删除整个模块", key=f"delmod_{mid}", type="secondary"):
                st.session_state[f"confirm_del_{mid}"] = True

            if st.session_state.get(f"confirm_del_{mid}"):
                st.warning(f"确认删除 [{mod['title']}]？此操作将清除该模块的所有学员进度和考试记录。")
                c1, c2 = st.columns(2)
                if c1.button("确认删除", key=f"yes_del_{mid}", type="primary"):
                    delete_module_db(mid)
                    st.session_state.pop(f"confirm_del_{mid}", None)
                    st.success("已删除")
                    st.rerun()
                if c2.button("取消", key=f"no_del_{mid}"):
                    st.session_state.pop(f"confirm_del_{mid}", None)
                    st.rerun()

            if qs:
                if st.button(f"清空题库 ({len(qs)} 题)", key=f"delexam_{mid}"):
                    save_exam_questions(mid, [])
                    st.success("题库已清空")
                    st.rerun()

# --------------- PAGE: STUDENT ANALYTICS ---------------
def page_analytics():
    if st.session_state.user.get("role") != "admin":
        st.warning("仅管理员可使用此功能")
        return
    st.markdown("## 考生分析")

    from db import get_supabase
    sb = get_supabase()
    result = sb.table("training_users").select("id, username, display_name, department, emp_id").eq("role", "user").order("id").execute()
    users = result.data

    if not users:
        st.info("暂无普通用户")
        return

    modules = get_modules()
    exams = get_exams()

    st.markdown("### 全员概览")
    header_cols = ["姓名", "部门", "工号"] + [mod["title"] for mod in modules.values()] + ["总进度"]
    rows = []
    for u in users:
        row = [u["display_name"], u["department"], u["emp_id"]]
        total_ch = 0
        done_ch = 0
        total_exams = 0
        passed_exams = 0
        for mid, mod in modules.items():
            ch_count = len(mod["chapters"])
            checks = get_read_checks(u["id"], mid, ch_count)
            read_n = sum(checks)
            best = get_best_exam(u["id"], mid)
            total_ch += ch_count
            done_ch += read_n
            exam_data = exams.get(mid, {})
            has_exam = bool(exam_data.get("questions") if isinstance(exam_data, dict) else exam_data)
            if has_exam:
                total_exams += 1
                if best is not None and best >= 80:
                    passed_exams += 1
            score_str = f"{best}分" if best is not None else "--"
            status = f"{read_n}/{ch_count}章 | {score_str}"
            row.append(status)
        total_items = total_ch + total_exams
        done_items = done_ch + passed_exams
        pct = round(done_items / total_items * 100) if total_items else 0
        row.append(f"{pct}%")
        rows.append(row)

    import pandas as pd
    df = pd.DataFrame(rows, columns=header_cols)
    st.dataframe(df, use_container_width=True, hide_index=True)

    st.markdown("---")
    st.markdown("### 个人详细分析")
    user_opts = {u["id"]: f"{u['display_name']} ({u['emp_id']})" for u in users}
    sel_uid = st.selectbox("选择考生", options=list(user_opts.keys()), format_func=lambda x: user_opts[x])

    sel_user = next(u for u in users if u["id"] == sel_uid)
    st.markdown(f"**{sel_user['display_name']}** | 部门: {sel_user['department']} | 工号: {sel_user['emp_id']}")

    for mid, mod in modules.items():
        st.markdown(f"#### {mod['title']}")
        ch_count = len(mod["chapters"])
        checks = get_read_checks(sel_uid, mid, ch_count)
        read_n = sum(checks)
        st.progress(read_n / ch_count if ch_count else 0, text=f"阅读进度: {read_n}/{ch_count}")

        exam_result = sb.table("training_exam_results").select("score, answers, taken_at").eq("user_id", sel_uid).eq("module_id", mid).order("taken_at", desc=True).execute()
        exam_rows = exam_result.data

        if exam_rows:
            best = max(r["score"] for r in exam_rows)
            st.markdown(f"考试次数: **{len(exam_rows)}** | 最高分: **{best}** | 最近一次: **{exam_rows[0]['score']}**")

            scores = [r["score"] for r in reversed(exam_rows)]
            if len(scores) > 1:
                st.line_chart({"得分": scores}, height=150)

            latest_ans = exam_rows[0]["answers"]
            try:
                ans_data = json.loads(latest_ans)
                details = ans_data.get("details", []) if isinstance(ans_data, dict) else []
            except (json.JSONDecodeError, TypeError):
                details = []

            # 需求3：考生个性化能力分析报告
            if details:
                wrong = [d for d in details if d["earned"] < d["max"]]
                if wrong:
                    st.markdown("#### 🔍 个性化能力诊断")
                    
                    error_types = {"single": 0, "multi": 0, "short": 0}
                    for d in wrong:
                        error_types[d.get("type", "single")] += 1
                        
                    col1, col2 = st.columns(2)
                    with col1:
                        st.warning(f"**核心薄弱点**：共 {len(wrong)} 处知识盲区")
                        st.markdown(f"- 概念记忆偏差（单/多选错误）: **{error_types['single'] + error_types['multi']}** 题")
                        st.markdown(f"- 流程与实操理解不足（简答题失分）: **{error_types['short']}** 题")
                        
                    with col2:
                        st.info("**💡 系统建议复习方向**")
                        if error_types['short'] > error_types['single'] + error_types['multi']:
                            st.markdown("- 建议重新阅读**系统流程与规范**章节，加强对底层逻辑的理解。")
                        else:
                            st.markdown("- 建议重点巩固**硬性指标与基础概念**，可利用碎片时间回顾模块大纲。")

                    with st.expander(f"薄弱环节明细与判分追踪 ({len(wrong)} 题)"):
                        for d in wrong:
                            type_label = {"single": "单选", "multi": "多选", "short": "简答"}.get(d.get("type", "single"), "")
                            st.markdown(f"**[{type_label}]** {d['q']} (得分 {d['earned']}/{d['max']})")
                            
                            if d.get("type") in ["single", "multi"]:
                                st.markdown(f"- *系统标准答案*：{d['correct']}")
                            else:
                                st.markdown(f"- *考生回答*：{d.get('user_ans', '未作答')}")
                                st.markdown(f"- *参考答案/关键词*：{d['correct']}")
                                if d.get("ai_reason"):
                                    # 管理员视角高亮展示 AI 判卷理由
                                    st.caption(f"**{d['ai_reason']}**")
                            st.divider()
                else:
                    st.success("🎉 最近一次考试表现完美，无知识盲区。")
        else:
            st.caption("尚未参加考试")

        t = get_training_time(sel_uid, mid)
        st.caption(f"学习时长: {fmt_time(t)}")

# --------------- SIDEBAR & MAIN ---------------
def main():
    # 公司名称从 Secrets 读取
    COMPANY_NAME = st.secrets.get("company_name", "培训系统")
    st.set_page_config(page_title=f"{COMPANY_NAME} - 培训系统", page_icon="📦", layout="wide")
    inject_global_css()

    # 初始化 Supabase 数据库
    init_db()

    # 种子管理员账号
    admin_pw = st.secrets.get("admin_password", None)
    if admin_pw:
        seed_admin(admin_pw)

    # 种子默认模块和考试（现在为空，由管理员上传）
    seed_default_modules(get_default_modules(), get_default_exams())

    if "page" not in st.session_state:
        st.session_state.page = "login"

    if "user" not in st.session_state:
        page_login()
        return

    user = st.session_state.user

    with st.sidebar:
        st.markdown(f"### 📦 {COMPANY_NAME}培训")
        st.markdown(f"**{user['display_name']}**")
        st.caption(f"👤 {user.get('department','')} | {user.get('emp_id','')}")
        st.markdown("---")

        if st.button("📊 培训仪表盘", use_container_width=True):
            st.session_state.page = "dashboard"
            st.rerun()

        if user["role"] == "admin":
            st.markdown("**管理员功能**")
            if st.button("👥 账号管理", use_container_width=True):
                st.session_state.page = "admin"
                st.rerun()

            if st.button("📁 上传课件", use_container_width=True):
                st.session_state.page = "upload_course"
                st.rerun()

            if st.button("📝 题库管理", use_container_width=True):
                st.session_state.page = "upload_exam"
                st.rerun()

            if st.button("⚙️ 课件管理", use_container_width=True):
                st.session_state.page = "manage_courses"
                st.rerun()

            if st.button("📈 考生分析", use_container_width=True):
                st.session_state.page = "analytics"
                st.rerun()

        st.markdown("---")
        if st.button("🚪 退出登录", use_container_width=True):
            for k in list(st.session_state.keys()):
                del st.session_state[k]
            st.rerun()

    page = st.session_state.page
    if page == "dashboard":
        page_dashboard()
    elif page == "admin" and user["role"] == "admin":
        page_admin()
    elif page == "module":
        page_module()
    elif page == "exam":
        page_exam()
    elif page == "result":
        page_result()
    elif page == "certificate":
        page_certificate()
    elif page == "upload_course":
        page_upload_course()
    elif page == "upload_exam":
        page_upload_exam()
    elif page == "manage_courses":
        page_manage_courses()
    elif page == "analytics":
        page_analytics()
    else:
        page_dashboard()

if __name__ == "__main__":
    main()